#!/usr/bin/env python3
# =============================================================================
# build_presentation.py — Gera a apresentação final (PPTX + PDF) da Árvore B,
# seguindo o modelo de slides do enunciado (Slides 1–10 + extras) e populando
# com os RESULTADOS REAIS dos experimentos (pastas results_*/ e _baseline/).
#
# Funde a intenção do guia APRESENTACAO_PPT.md (btree-main) com os dados reais
# desta implementação. Saída em trabalho/apresentacao/.
#
# Requer (venv): python-pptx, reportlab, matplotlib.
# Uso: python3 experiments/build_presentation.py
# =============================================================================
import csv
import os
from collections import defaultdict

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "apresentacao")
ASSETS = os.path.join(OUT, "assets")
GRAPHS = os.path.join(ROOT, "results_graphs")
os.makedirs(ASSETS, exist_ok=True)

ACCENT = "#1d4ed8"
ACCENT2 = "#dc2626"
INK = "#111827"
MUTED = "#6b7280"

FIG = (8.8, 4.9)   # tamanho dos gráficos (maior = mais legível)
DPI = 165


# --------------------------------------------------------------------------- #
# Carregamento de dados
# --------------------------------------------------------------------------- #
def load(path):
    rows = []
    if not os.path.exists(path):
        return rows
    with open(path, newline="") as f:
        for r in csv.DictReader(f):
            for k in ("m", "N", "disk_accesses", "height", "total_nodes",
                      "free_nodes", "file_bytes"):
                if r.get(k, "") != "":
                    r[k] = int(float(r[k]))
            for k in ("avg_io_per_op", "wall_ms", "cpu_user_ms",
                      "cpu_sys_ms", "io_wait_ms"):
                if r.get(k, "") != "":
                    r[k] = float(r[k])
            rows.append(r)
    return rows


ORDER = load(os.path.join(ROOT, "results_order_m_impact", "data.csv"))
SCALE = load(os.path.join(ROOT, "results_set_size_scaling", "data.csv"))
REUSE = load(os.path.join(ROOT, "results_node_reuse", "data.csv"))
ORDER_LAP = load(os.path.join(ROOT, "_baseline", "results_order_m_impact", "data.csv"))


def by_phase(rows, phase, mode=None):
    return [r for r in rows if r["phase"] == phase and (not mode or r["mode"] == mode)]


# --------------------------------------------------------------------------- #
# Gráficos (matplotlib -> PNG)
# --------------------------------------------------------------------------- #
def _save(fig, name):
    p = os.path.join(ASSETS, name)
    fig.tight_layout(); fig.savefig(p); plt.close(fig)
    return p


def chart_io_vs_m():
    fig, ax = plt.subplots(figsize=FIG, dpi=DPI)
    for mode, color, lbl in (("rand", ACCENT, "aleatório"), ("seq", ACCENT2, "sequencial")):
        pts = sorted((r["m"], r["avg_io_per_op"]) for r in by_phase(ORDER, "search", mode))
        if pts:
            xs, ys = zip(*pts); ax.plot(xs, ys, "o-", color=color, label=lbl, linewidth=2.4, markersize=6)
    ax.set_xscale("log"); ax.set_xlabel("ordem m (escala log)"); ax.set_ylabel("I/O médio por busca")
    ax.set_title("Acessos a disco por busca vs ordem m  (N = 100.000)")
    ax.grid(True, alpha=0.3); ax.legend(fontsize=11)
    return _save(fig, "chart_io_vs_m.png")


def chart_height_vs_m():
    fig, ax = plt.subplots(figsize=FIG, dpi=DPI)
    for mode, color, lbl in (("rand", ACCENT, "aleatório"), ("seq", ACCENT2, "sequencial")):
        pts = sorted((r["m"], r["height"]) for r in by_phase(ORDER, "insert", mode))
        if pts:
            xs, ys = zip(*pts); ax.plot(xs, ys, "s-", color=color, label=lbl, linewidth=2.4, markersize=6)
    ax.set_xscale("log"); ax.set_xlabel("ordem m (escala log)"); ax.set_ylabel("altura da árvore (níveis)")
    ax.set_title("Altura da árvore vs ordem m  (N = 100.000)")
    ax.grid(True, alpha=0.3); ax.legend(fontsize=11)
    return _save(fig, "chart_height_vs_m.png")


def chart_io_vs_N():
    fig, ax = plt.subplots(figsize=FIG, dpi=DPI)
    colors = {3: ACCENT, 100: "#059669", 1000: ACCENT2}
    for m in (3, 100, 1000):
        pts = sorted((r["N"], r["avg_io_per_op"]) for r in by_phase(SCALE, "search", "rand") if r["m"] == m)
        if pts:
            xs, ys = zip(*pts); ax.plot(xs, ys, "o-", color=colors[m], label=f"m = {m}", linewidth=2.4, markersize=6)
    ax.set_xscale("log"); ax.set_xlabel("N — nº de chaves (escala log)"); ax.set_ylabel("I/O médio por busca")
    ax.set_title("Escalabilidade: I/O por busca vs N  (aleatório)")
    ax.grid(True, alpha=0.3); ax.legend(fontsize=11)
    return _save(fig, "chart_io_vs_N.png")


def chart_reuse():
    import numpy as np
    ref = [r for r in REUSE if r["phase"] == "churn_refill"]
    by = defaultdict(dict)
    for r in ref:
        by[(r["m"], r["N"])][r["reuse"]] = r
    keys = sorted(by)
    labels = [f"m={m}\nN={N:,}" for (m, N) in keys]
    off = [by[k].get("0", {}).get("file_bytes", 0) / 1024 for k in keys]
    on = [by[k].get("1", {}).get("file_bytes", 0) / 1024 for k in keys]
    x = np.arange(len(keys)); w = 0.38
    fig, ax = plt.subplots(figsize=FIG, dpi=DPI)
    ax.bar(x - w/2, off, w, label="sem reaproveitamento", color=ACCENT2)
    ax.bar(x + w/2, on, w, label="com reaproveitamento", color=ACCENT)
    ax.set_xticks(x); ax.set_xticklabels(labels, fontsize=9)
    ax.set_ylabel("tamanho do arquivo (KB)")
    ax.set_title("Ocupação do arquivo após churn: com vs sem free list")
    ax.grid(True, axis="y", alpha=0.3); ax.legend(fontsize=11)
    return _save(fig, "chart_reuse.png")


def chart_machines():
    import numpy as np

    def total_by_m(rows):
        tot = defaultdict(float)
        for r in rows:
            if r["mode"] == "rand" and r["phase"] in ("insert", "search", "delete"):
                tot[r["m"]] += r["wall_ms"]
        return tot
    lap, tit = total_by_m(ORDER_LAP), total_by_m(ORDER)
    ms = sorted(set(lap) & set(tit))
    if not ms:
        return None
    x = np.arange(len(ms)); w = 0.38
    fig, ax = plt.subplots(figsize=FIG, dpi=DPI)
    ax.bar(x - w/2, [lap[m]/1000 for m in ms], w, label="Notebook", color=ACCENT)
    ax.bar(x + w/2, [tit[m]/1000 for m in ms], w, label="Titan (USP)", color="#7c3aed")
    ax.set_xticks(x); ax.set_xticklabels([str(m) for m in ms], fontsize=9)
    ax.set_xlabel("ordem m"); ax.set_ylabel("tempo total ins+busca+rem (s)")
    ax.set_title("Tempo por máquina (N=100k, aleatório) — acessos a disco idênticos")
    ax.grid(True, axis="y", alpha=0.3); ax.legend(fontsize=11)
    return _save(fig, "chart_machines.png")


print("gerando gráficos...")
C_IO_M = chart_io_vs_m()
C_H_M = chart_height_vs_m()
C_IO_N = chart_io_vs_N()
C_REUSE = chart_reuse()
C_MACH = chart_machines()


# --------------------------------------------------------------------------- #
# Tabelas dinâmicas (dados reais)
# --------------------------------------------------------------------------- #
def order_table():
    by_m = defaultdict(dict)
    for r in ORDER:
        if r["mode"] == "rand":
            by_m[r["m"]][r["phase"]] = r
    headers = ["m", "altura", "busca I/O/op", "inser. I/O/op", "rem. I/O/op", "arquivo (KB)"]
    rows = []
    for m in sorted(by_m):
        d = by_m[m]
        rows.append([str(m), str(d["insert"]["height"]),
                     f'{d["search"]["avg_io_per_op"]:.2f}',
                     f'{d["insert"]["avg_io_per_op"]:.2f}',
                     f'{d["delete"]["avg_io_per_op"]:.2f}',
                     f'{d["insert"]["file_bytes"]/1024:,.0f}'])
    return headers, rows


def reuse_table():
    ref = [r for r in REUSE if r["phase"] == "churn_refill"]
    by = defaultdict(dict)
    for r in ref:
        by[(r["m"], r["N"])][r["reuse"]] = r
    headers = ["m", "N", "sem reuso (KB)", "com reuso (KB)", "economia"]
    rows = []
    for (m, N) in sorted(by):
        on = by[(m, N)].get("1"); off = by[(m, N)].get("0")
        if on and off and off["file_bytes"]:
            save = (1 - on["file_bytes"]/off["file_bytes"]) * 100
            rows.append([str(m), f"{N:,}", f'{off["file_bytes"]/1024:,.0f}',
                         f'{on["file_bytes"]/1024:,.0f}', f"{save:.0f}%"])
    return headers, rows


ORDER_HEADERS, ORDER_ROWS = order_table()
REUSE_HEADERS, REUSE_ROWS = reuse_table()


# --------------------------------------------------------------------------- #
# Modelo de slides
# --------------------------------------------------------------------------- #
def G(name):
    return os.path.join(GRAPHS, name)

SLIDES = [
    {   # 1 — capa
        "layout": "cover",
        "title": "Árvore B em Memória Secundária",
        "subtitle": "Implementação de uma classe Árvore B de ordem m operando estritamente em disco",
        "bullets": [
            "Algoritmos e Estruturas de Dados — AED-PG-2026 (5955001-3)",
            "Prof. Dr. José Augusto Baranauskas — DCM / USP",
            "Linguagem: C++17  ·  1º Semestre/2026",
            "Alunos: Cássio Takarada  ·  Cézio Luiz Ferreira Junior",
        ],
    },
    {   # 2 — decisões (imagem larga em faixa)
        "layout": "image_full",
        "title": "2. Decisões de implementação",
        "bullets": [
            "Ordem m é constante simbólica de compilação (M): estrutura totalmente parametrizada.",
            "Raiz da árvore: armazenada no header (registro 0) do arquivo — campo root.",
            "Um nó por I/O: readNode/writeNode de um registro de PAGE_SIZE; a árvore NUNCA é carregada inteira em memória.",
            "Registro do nó: n  ·  A[0..m-1] (ponteiros RRN)  ·  K[0..m-1] (chaves).",
            "Reaproveitamento de nós: estrutura = lista encadeada de livres (pilha LIFO) no próprio arquivo - free_head no header; cada nó livre marca n=-1 e aponta o próximo em K[1]; allocNode reusa antes de estender o arquivo.",
        ],
        "images": [(G("m3_final.png"), "Árvore B (m=3) com 40 chaves — exportada do nosso código (Graphviz); cada caixa mostra o RRN do nó em disco")],
    },
    {   # 3 — métodos
        "layout": "bullets",
        "title": "3. Métodos implementados",
        "bullets": [
            "mSearch / mSearchPath - busca m-way da raiz à folha; retorna (RRN, índice, achou).",
            "insertB - inserção bottom-up com propagação de split (Equação 1).",
            "deleteB - remoção com sucessor in-order, redistribuição (empréstimo) e fusão (merge).",
            "Auxiliares: splitNode, insertInNode, removeFromNode, findSuccessorLeaf.",
            "Monitoramento: printTree (hierárquico), height, exportDot (Graphviz).",
            "DiskManager: readNode, writeNode, allocNode, freeNode, contador de acessos, setReuse.",
            "bench - driver experimental não interativo (CSV de métricas).",
        ],
    },
    {   # 3b — RESUMO de tudo que o código faz (núcleo + extras)
        "layout": "two_col",
        "title": "Resumo — tudo que o código faz (núcleo + extras)",
        "left_title": "Núcleo — Árvore B em disco",
        "left": [
            "Árvore B de ordem m parametrizável (M em tempo de compilação).",
            "Opera 100% em disco: 1 nó por I/O; nada é carregado em RAM.",
            "Busca m-way (mSearch / mSearchPath).",
            "Inserção bottom-up com split (insertB).",
            "Remoção com sucessor, redistribuição e fusão (deleteB).",
            "Persistência: header (root, total, free_head) + nós de tamanho fixo, acesso direto por RRN.",
        ],
        "right_title": "Extras / funcionalidades acessórias",
        "right": [
            "Contador de acessos ao disco (métrica central da avaliação).",
            "Reaproveitamento de nós via free list, com liga/desliga (setReuse).",
            "Impressão hierárquica (printTree) e cálculo de altura (height).",
            "Exportação Graphviz (exportDot) -> imagens PNG dos grafos.",
            "Medição de ocupação: total_nodes, free_nodes, fileSizeBytes.",
            "Benchmark (bench): CSV com I/O, tempo CPU/IO (getrusage), altura, nós, tamanho.",
            "Automação: run_all.sh + make_tables.py (tabelas/gráficos) + compare.py (2 máquinas).",
            "Testes de stress + menu interativo (CRUD + métricas).",
        ],
    },
    {   # 4 — experimentos
        "layout": "text_image",
        "title": "4. Análises efetuadas (experimentos)",
        "bullets": [
            "Exp. 1 - Impacto da ordem m: m em {3,4,5,8,16,32,64,100,128,256,512,1000}, N=10^5.",
            "Exp. 2 - Escala do conjunto: N em {10^3,10^4,10^5,10^6}, m em {3,100,1000}.",
            "Exp. 3 - Ocupação do arquivo: churn COM e SEM reaproveitamento.",
            "Exp. 4 - Tempo: CPU (user+sys) vs espera de I/O (getrusage).",
            "Modos aleatório e sequencial em todos.",
            "Novidade: validação em 2 máquinas (Notebook x Titan/USP) - disco idêntico.",
        ],
        "images": [(C_IO_N, None)],
    },
    {   # 5 — métricas
        "layout": "text_image",
        "title": "5. Métricas utilizadas",
        "bullets": [
            "Acessos ao disco - contador readNode+writeNode; total e média/op (métrica honesta de I/O).",
            "Altura da árvore - nº de níveis (~ log_m N).",
            "Ocupação - nós físicos, nós livres e tamanho em bytes.",
            "Tempo de parede (wall) - relógio de alta resolução.",
            "Tempo de CPU - usuário (lógica) e sistema (syscalls de I/O), via getrusage.",
            "Espera de I/O - wall menos CPU.",
        ],
        "images": [(C_H_M, None)],
    },
    {   # 6 — tabela impacto de m
        "layout": "table_image",
        "title": "6. Tabela de Resultados - impacto de m (N=10^5, aleatório)",
        "table": (ORDER_HEADERS, ORDER_ROWS),
        "images": [(C_IO_M, None)],
        "footnote": "I/O por busca cai de 13,25 (m=3, altura 14) para 2,00 (m>=512, altura 2). Ganho satura por volta de m~64-128.",
    },
    {   # 7 — reaproveitamento
        "layout": "table_image",
        "title": "6b. Resultados - reaproveitamento de nós (churn)",
        "table": (REUSE_HEADERS, REUSE_ROWS),
        "images": [(C_REUSE, None)],
        "footnote": "O reaproveitamento de nós economiza ~27-30% do tamanho do arquivo neste workload.",
    },
    {   # 8 — LLM + máquinas
        "layout": "text_image",
        "title": "7. Utilização de LLM",
        "bullets": [
            "Ferramenta: Claude (Anthropic) via Claude Code CLI - auxílio ao desenvolvimento.",
            "Síntese e discussão dos pseudocódigos dos slides do professor.",
            "Identificação e correção de 2 bugs: (1) stale-header em insertB; (2) eofbit do fstream.",
            "Geração do harness experimental, tabelas, gráficos e desta apresentação.",
            "Todo o código foi revisado e validado manualmente pelo autor.",
        ],
        "images": [(C_MACH, None)] if C_MACH else [],
    },
    {   # 9 — dificuldades
        "layout": "bullets",
        "title": "8. Dificuldades · Vantagens x Desvantagens",
        "bullets": [
            "Dificuldades: indexação 1-based dos nós; ordem do path no reparo de underflow; eofbit do fstream; garantir 1 nó por I/O (nada em RAM).",
            "Vantagens da Árvore B: altura baixa => pouquíssimos acessos a disco; sempre balanceada; ideal para memória secundária / SGBDs.",
            "Desvantagens: nós podem ficar subocupados (~50% no caso sequencial); remoção é complexa; para m muito grande, a busca dentro do nó passa a custar CPU.",
            "Trade-off central: m maior => menos I/O, mais CPU por nó.",
        ],
    },
    {   # 9 — aplicações práticas
        "layout": "bullets",
        "title": "9. Aplicações práticas",
        "bullets": [
            "Índices de SGBDs: B-tree / B+ tree são a base dos índices de MySQL/InnoDB, PostgreSQL, Oracle e SQL Server.",
            "Sistemas de arquivos: NTFS, HFS+/APFS, ext4 (HTree) e Btrfs usam B-trees para diretórios e metadados.",
            "Armazenamento chave-valor / embarcado: BerkeleyDB, SQLite, LMDB.",
            "Onde eu usaria: como índice (primário ou secundário) de um banco de dados em disco - exatamente o cenário deste trabalho.",
        ],
    },
    {   # 10 — conclusão e referências
        "layout": "bullets",
        "title": "10. Conclusão e Referências",
        "bullets": [
            "O I/O por operação cai com m até SATURAR (ponto ~ m=64-128): existe um m ótimo (nó dimensionado para um bloco de disco).",
            "O reaproveitamento de nós (free list) economiza ~27-30% do tamanho do arquivo no workload de churn.",
            "Resultados determinísticos: acessos a disco IDÊNTICOS em 2 máquinas (Notebook x Titan/USP); só o tempo varia.",
            "Referências:",
            "   - Comer, D. (1979). The Ubiquitous B-Tree. ACM Computing Surveys.",
            "   - Knuth, D. The Art of Computer Programming, Vol. 3 - Sorting and Searching.",
            "   - Folk, M. & Zoellick, B. File Structures.",
            "   - Baranauskas, J. A. Slides AED-PG-2026 (Árvores B).",
        ],
    },
]


# --------------------------------------------------------------------------- #
# PPTX
# --------------------------------------------------------------------------- #
def build_pptx(path):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from PIL import Image

    def rgb(h):
        return RGBColor.from_string(h.lstrip("#"))

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    SW = prs.slide_width
    blank = prs.slide_layouts[6]

    def add_text(slide, left, top, width, height, runs, anchor=MSO_ANCHOR.TOP, space=6):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        first = True
        for (text, size, color, bold, bullet) in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False; p.space_after = Pt(space)
            r = p.add_run(); r.text = ("•  " + text) if bullet else text
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = rgb(color); r.font.name = "Calibri"
        return tb

    def add_table(slide, headers, rows, left, top, width, height):
        gtbl = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
        for j, h in enumerate(headers):
            c = gtbl.cell(0, j); c.text = h
            pr = c.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
            run = pr.runs[0]; run.font.size = Pt(11); run.font.bold = True
            run.font.color.rgb = rgb("#ffffff")
            c.fill.solid(); c.fill.fore_color.rgb = rgb(ACCENT)
        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                c = gtbl.cell(i, j); c.text = str(val)
                pr = c.text_frame.paragraphs[0]
                pr.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                run = pr.runs[0]; run.font.size = Pt(10); run.font.color.rgb = rgb(INK)
                c.fill.solid(); c.fill.fore_color.rgb = rgb("#eef2ff" if i % 2 == 0 else "#ffffff")
        return gtbl

    def add_image_fit(slide, img, left, top, max_w, max_h, caption=None):
        try:
            iw, ih = Image.open(img).size
        except Exception:
            iw, ih = (1280, 720)
        scale = min(max_w / iw, max_h / ih)
        w, h = int(iw * scale), int(ih * scale)
        l = left + (max_w - w) // 2
        slide.shapes.add_picture(img, l, top, width=w, height=h)
        if caption:
            add_text(slide, left, top + h + Emu(30000), max_w, Inches(0.4),
                     [(caption, 9, MUTED, False, False)])

    def title_band(slide, text):
        band = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.95))
        band.fill.solid(); band.fill.fore_color.rgb = rgb(ACCENT); band.line.fill.background()
        add_text(slide, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.72),
                 [(text, 23, "#ffffff", True, False)], anchor=MSO_ANCHOR.MIDDLE)

    def footer(slide, idx):
        add_text(slide, Inches(0.5), Inches(7.06), Inches(12.3), Inches(0.34),
                 [(f"Árvore B em Memória Secundária — AED-PG-2026          {idx+1}/{len(SLIDES)}",
                   9, MUTED, False, False)])

    for idx, s in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank)
        lay = s["layout"]

        if lay == "cover":
            band = slide.shapes.add_shape(1, 0, Inches(2.2), SW, Inches(1.7))
            band.fill.solid(); band.fill.fore_color.rgb = rgb(ACCENT); band.line.fill.background()
            add_text(slide, Inches(0.8), Inches(2.35), Inches(11.7), Inches(1.4),
                     [(s["title"], 40, "#ffffff", True, False)], anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.9),
                     [(s["subtitle"], 18, INK, False, False)])
            add_text(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.0),
                     [(b, 15, MUTED, False, False) for b in s["bullets"]])
            continue

        title_band(slide, s["title"])
        body_top = Inches(1.15)

        if lay == "image_full":
            add_text(slide, Inches(0.5), body_top, Inches(12.3), Inches(2.95),
                     [(b, 13, INK, False, True) for b in s["bullets"]], space=5)
            img, cap = s["images"][0]
            add_image_fit(slide, img, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.65), cap)

        elif lay == "two_col":
            lruns = [(s["left_title"], 15, ACCENT, True, False)] + \
                    [(b, 12.5, INK, False, True) for b in s["left"]]
            rruns = [(s["right_title"], 15, ACCENT, True, False)] + \
                    [(b, 12.5, INK, False, True) for b in s["right"]]
            add_text(slide, Inches(0.45), body_top, Inches(6.15), Inches(5.9), lruns, space=5)
            add_text(slide, Inches(6.85), body_top, Inches(6.15), Inches(5.9), rruns, space=5)

        elif lay == "text_image":
            add_text(slide, Inches(0.5), body_top, Inches(6.4), Inches(5.8),
                     [(b, 14, INK, False, True) for b in s["bullets"]])
            img, cap = s["images"][0]
            add_image_fit(slide, img, Inches(6.95), Inches(1.5), Inches(6.1), Inches(5.0), cap)

        elif lay == "table_image":
            headers, rows = s["table"]
            add_table(slide, headers, rows, Inches(0.4), body_top, Inches(5.9), Inches(5.0))
            img, cap = s["images"][0]
            add_image_fit(slide, img, Inches(6.55), Inches(1.4), Inches(6.5), Inches(5.1), cap)
            if s.get("footnote"):
                add_text(slide, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.5),
                         [(s["footnote"], 11, MUTED, False, False)])

        else:  # bullets
            add_text(slide, Inches(0.7), body_top, Inches(12.0), Inches(5.8),
                     [(b, 16, INK, False, True) for b in s["bullets"]])

        footer(slide, idx)

    prs.save(path)
    print("PPTX:", path)


# --------------------------------------------------------------------------- #
# PDF (reportlab)
# --------------------------------------------------------------------------- #
def build_pdf(path):
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    from reportlab.lib.colors import HexColor, white

    W, H = 960, 540
    c = canvas.Canvas(path, pagesize=(W, H))

    def wrap(text, font, size, max_w):
        words = text.split(" "); lines, cur = [], ""
        for w in words:
            t = (cur + " " + w).strip()
            if c.stringWidth(t, font, size) <= max_w:
                cur = t
            else:
                if cur:
                    lines.append(cur)
                cur = w
        if cur:
            lines.append(cur)
        return lines

    def bullets(items, x, y, max_w, size=14, lh=20, gap=6):
        for b in items:
            lines = wrap(b, "Helvetica", size, max_w)
            c.setFillColor(HexColor(ACCENT)); c.setFont("Helvetica", size)
            c.drawString(x, y, "•")
            c.setFillColor(HexColor(INK))
            for ln in lines:
                c.drawString(x + 16, y, ln); y -= lh
            y -= gap
        return y

    def draw_image_fit(img, x, y_top, max_w, max_h, caption=None):
        try:
            ir = ImageReader(img); iw, ih = ir.getSize()
        except Exception:
            return
        scale = min(max_w / iw, max_h / ih)
        w, h = iw * scale, ih * scale
        c.drawImage(ir, x + (max_w - w) / 2, y_top - h, width=w, height=h,
                    preserveAspectRatio=True, mask="auto")
        if caption:
            c.setFont("Helvetica", 8.5); c.setFillColor(HexColor(MUTED))
            for k, ln in enumerate(wrap(caption, "Helvetica", 8.5, max_w)):
                c.drawCentredString(x + max_w / 2, y_top - h - 12 - k * 11, ln)

    def draw_table(headers, rows, x, y_top, total_w):
        nc = len(headers); cw = total_w / nc; rh = 21
        c.setFillColor(HexColor(ACCENT)); c.rect(x, y_top - rh, total_w, rh, fill=1, stroke=0)
        c.setFillColor(white); c.setFont("Helvetica-Bold", 9)
        for j, hd in enumerate(headers):
            c.drawCentredString(x + cw * j + cw / 2, y_top - rh + 7, str(hd))
        c.setFont("Helvetica", 9); yy = y_top - rh
        for i, row in enumerate(rows):
            if i % 2 == 0:
                c.setFillColor(HexColor("#eef2ff")); c.rect(x, yy - rh, total_w, rh, fill=1, stroke=0)
            c.setFillColor(HexColor(INK))
            for j, val in enumerate(row):
                if j == 0:
                    c.drawString(x + 6, yy - rh + 7, str(val))
                else:
                    c.drawCentredString(x + cw * j + cw / 2, yy - rh + 7, str(val))
            yy -= rh
        return yy

    for idx, s in enumerate(SLIDES):
        lay = s["layout"]
        c.setFillColor(HexColor("#ffffff")); c.rect(0, 0, W, H, fill=1, stroke=0)

        if lay == "cover":
            c.setFillColor(HexColor(ACCENT)); c.rect(0, H - 250, W, 130, fill=1, stroke=0)
            c.setFillColor(white); c.setFont("Helvetica-Bold", 34)
            for k, ln in enumerate(wrap(s["title"], "Helvetica-Bold", 34, W - 120)):
                c.drawString(60, H - 175 - k * 38, ln)
            c.setFillColor(HexColor(INK)); c.setFont("Helvetica", 15)
            for k, ln in enumerate(wrap(s["subtitle"], "Helvetica", 15, W - 120)):
                c.drawString(60, H - 285 - k * 20, ln)
            c.setFillColor(HexColor(MUTED)); c.setFont("Helvetica", 13)
            for k, b in enumerate(s["bullets"]):
                c.drawString(60, H - 360 - k * 24, b)
            c.showPage(); continue

        c.setFillColor(HexColor(ACCENT)); c.rect(0, H - 60, W, 60, fill=1, stroke=0)
        c.setFillColor(white); c.setFont("Helvetica-Bold", 18)
        c.drawString(36, H - 40, s["title"])
        top = H - 88

        if lay == "image_full":
            bullets(s["bullets"], 36, top, W - 90, size=12.5, lh=17, gap=5)
            img, cap = s["images"][0]
            draw_image_fit(img, 30, 220, W - 60, 180, cap)

        elif lay == "two_col":
            c.setFillColor(HexColor(ACCENT)); c.setFont("Helvetica-Bold", 13)
            c.drawString(36, top, s["left_title"]); c.drawString(W/2 + 10, top, s["right_title"])
            bullets(s["left"], 36, top - 24, W/2 - 70, size=10.5, lh=14, gap=5)
            bullets(s["right"], W/2 + 10, top - 24, W/2 - 60, size=10.5, lh=14, gap=5)

        elif lay == "text_image":
            bullets(s["bullets"], 36, top, 470, size=12.5, lh=17, gap=5)
            img, cap = s["images"][0]
            draw_image_fit(img, 500, top + 6, 430, 380, cap)

        elif lay == "table_image":
            headers, rows = s["table"]
            draw_table(headers, rows, 28, top, 430)
            img, cap = s["images"][0]
            draw_image_fit(img, 470, top, 470, 380, cap)
            if s.get("footnote"):
                c.setFillColor(HexColor(MUTED)); c.setFont("Helvetica-Oblique", 10)
                for k, ln in enumerate(wrap(s["footnote"], "Helvetica-Oblique", 10, W - 70)):
                    c.drawString(36, 56 - k * 13, ln)

        else:  # bullets
            bullets(s["bullets"], 40, top, W - 110, size=14, lh=20, gap=6)

        c.setFillColor(HexColor(MUTED)); c.setFont("Helvetica", 8)
        c.drawString(36, 16, "Árvore B em Memória Secundária — AED-PG-2026")
        c.drawRightString(W - 36, 16, f"{idx+1}/{len(SLIDES)}")
        c.showPage()

    c.save()
    print("PDF:", path)


if __name__ == "__main__":
    build_pptx(os.path.join(OUT, "APRESENTACAO.pptx"))
    build_pdf(os.path.join(OUT, "APRESENTACAO.pdf"))
    print("\nApresentação gerada em:", OUT)
