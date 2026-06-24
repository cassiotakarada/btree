#!/usr/bin/env python3
import csv
import os
from collections import defaultdict

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "apresentacao")
ASSETS = os.path.join(OUT, "assets")
GRAPHS = os.path.join(ROOT, "results_graphs")
os.makedirs(ASSETS, exist_ok=True)

ACCENT = "#1d4ed8"
ACCENT2 = "#dc2626"
INK = "#111827"
MUTED = "#6b7280"

FIG = (8.8, 4.9)
DPI = 165

def load(path):
    rows = []
    if not os.path.exists(path):
        return rows
    with open(path, newline="") as f:
        for r in csv.DictReader(f):
            for k in ("m", "N", "disk_accesses", "height", "total_nodes",
                      "free_nodes", "file_bytes"):
                if r.get(k, "") != "":
                    r[k] = int(float(r[k]))
            for k in ("avg_io_per_op", "wall_ms", "cpu_user_ms",
                      "cpu_sys_ms", "io_wait_ms"):
                if r.get(k, "") != "":
                    r[k] = float(r[k])
            rows.append(r)
    return rows

ORDER = load(os.path.join(ROOT, "results_order_m_impact", "data.csv"))
SCALE = load(os.path.join(ROOT, "results_set_size_scaling", "data.csv"))
REUSE = load(os.path.join(ROOT, "results_node_reuse", "data.csv"))
ORDER_LAP = load(os.path.join(ROOT, "_baseline", "results_order_m_impact", "data.csv"))
SCALE_LAP = load(os.path.join(ROOT, "_baseline", "results_set_size_scaling", "data.csv"))

TITAN = "#7c3aed"  # cor do servidor Titan (USP)

def by_phase(rows, phase, mode=None):
    return [r for r in rows if r["phase"] == phase and (not mode or r["mode"] == mode)]

def _save(fig, name):
    p = os.path.join(ASSETS, name)
    fig.tight_layout(); fig.savefig(p); plt.close(fig)
    return p

def chart_io_vs_m():
    fig, ax = plt.subplots(figsize=FIG, dpi=DPI)
    for mode, color, lbl in (("rand", ACCENT, "aleatório"), ("seq", ACCENT2, "sequencial")):
        pts = sorted((r["m"], r["avg_io_per_op"]) for r in by_phase(ORDER, "search", mode))
        if pts:
            xs, ys = zip(*pts); ax.plot(xs, ys, "o-", color=color, label=lbl, linewidth=2.4, markersize=6)
    ax.set_xscale("log"); ax.set_xlabel("ordem m (escala log)"); ax.set_ylabel("I/O médio por busca")
    ax.set_title("Acessos a disco por busca vs ordem m  (N = 100.000)")
    ax.grid(True, alpha=0.3); ax.legend(fontsize=11)
    return _save(fig, "chart_io_vs_m.png")

def chart_height_vs_m():
    fig, ax = plt.subplots(figsize=FIG, dpi=DPI)
    for mode, color, lbl in (("rand", ACCENT, "aleatório"), ("seq", ACCENT2, "sequencial")):
        pts = sorted((r["m"], r["height"]) for r in by_phase(ORDER, "insert", mode))
        if pts:
            xs, ys = zip(*pts); ax.plot(xs, ys, "s-", color=color, label=lbl, linewidth=2.4, markersize=6)
    ax.set_xscale("log"); ax.set_xlabel("ordem m (escala log)"); ax.set_ylabel("altura da árvore (níveis)")
    ax.set_title("Altura da árvore vs ordem m  (N = 100.000)")
    ax.grid(True, alpha=0.3); ax.legend(fontsize=11)
    return _save(fig, "chart_height_vs_m.png")

def chart_io_vs_N():
    fig, ax = plt.subplots(figsize=FIG, dpi=DPI)
    colors = {3: ACCENT, 100: "#059669", 1000: ACCENT2}
    for m in (3, 100, 1000):
        pts = sorted((r["N"], r["avg_io_per_op"]) for r in by_phase(SCALE, "search", "rand") if r["m"] == m)
        if pts:
            xs, ys = zip(*pts); ax.plot(xs, ys, "o-", color=colors[m], label=f"m = {m}", linewidth=2.4, markersize=6)
    ax.set_xscale("log"); ax.set_xlabel("N — nº de chaves (escala log)"); ax.set_ylabel("I/O médio por busca")
    ax.set_title("Escalabilidade: I/O por busca vs N  (aleatório)")
    ax.grid(True, alpha=0.3); ax.legend(fontsize=11)
    return _save(fig, "chart_io_vs_N.png")

def chart_reuse():
    import numpy as np
    ref = [r for r in REUSE if r["phase"] == "churn_refill"]
    by = defaultdict(dict)
    for r in ref:
        by[(r["m"], r["N"])][r["reuse"]] = r
    keys = sorted(by)
    labels = [f"m={m}\nN={N:,}" for (m, N) in keys]
    off = [by[k].get("0", {}).get("file_bytes", 0) / 1024 for k in keys]
    on = [by[k].get("1", {}).get("file_bytes", 0) / 1024 for k in keys]
    x = np.arange(len(keys)); w = 0.38
    fig, ax = plt.subplots(figsize=FIG, dpi=DPI)
    ax.bar(x - w/2, off, w, label="sem reaproveitamento", color=ACCENT2)
    ax.bar(x + w/2, on, w, label="com reaproveitamento", color=ACCENT)
    ax.set_xticks(x); ax.set_xticklabels(labels, fontsize=9)
    ax.set_ylabel("tamanho do arquivo (KB)")
    ax.set_title("Ocupação do arquivo após churn: com vs sem free list")
    ax.grid(True, axis="y", alpha=0.3); ax.legend(fontsize=11)
    return _save(fig, "chart_reuse.png")

def chart_machines():
    import numpy as np

    def total_by_m(rows, metric):
        tot = defaultdict(float)
        for r in rows:
            if r["mode"] == "rand" and r["phase"] in ("insert", "search", "delete"):
                tot[r["m"]] += metric(r)
        return tot

    panels = [
        ("Wall time (relógio)", lambda r: r["wall_ms"]),
        ("Tempo de CPU (user+sys)", lambda r: r["cpu_user_ms"] + r["cpu_sys_ms"]),
    ]
    fig, axes = plt.subplots(1, 2, figsize=FIG, dpi=DPI, sharey=True)
    drew = False
    for ax, (title, metric) in zip(axes, panels):
        lap, tit = total_by_m(ORDER_LAP, metric), total_by_m(ORDER, metric)
        ms = sorted(set(lap) & set(tit))
        if not ms:
            continue
        drew = True
        x = np.arange(len(ms)); w = 0.38
        ax.bar(x - w/2, [lap[m]/1000 for m in ms], w, label="Notebook", color=ACCENT)
        ax.bar(x + w/2, [tit[m]/1000 for m in ms], w, label="Titan (USP)", color="#7c3aed")
        ax.set_xticks(x); ax.set_xticklabels([str(m) for m in ms], fontsize=8)
        ax.set_xlabel("ordem m"); ax.set_title(title, fontsize=11)
        ax.grid(True, axis="y", alpha=0.3)
    if not drew:
        plt.close(fig)
        return None
    axes[0].set_ylabel("tempo total ins+busca+rem (s)")
    axes[0].legend(fontsize=10)
    fig.suptitle("Tempo por máquina (N=100k, aleatório) — acessos a disco idênticos", fontsize=12)
    return _save(fig, "chart_machines.png")

def chart_phase_io():
    """Resultados SEPARADOS por operação: I/O médio de inserção/busca/remoção."""
    import numpy as np
    ms = [3, 8, 32, 100, 1000]
    by = defaultdict(dict)
    for r in ORDER:
        if r["mode"] == "rand":
            by[r["m"]][r["phase"]] = r
    phases = [("insert", "inserção", ACCENT),
              ("search", "busca", "#059669"),
              ("delete", "remoção", ACCENT2)]
    x = np.arange(len(ms)); w = 0.26
    fig, ax = plt.subplots(figsize=FIG, dpi=DPI)
    for i, (ph, lbl, col) in enumerate(phases):
        ys = [by[m].get(ph, {}).get("avg_io_per_op", 0) for m in ms]
        ax.bar(x + (i - 1) * w, ys, w, label=lbl, color=col)
    ax.set_xticks(x); ax.set_xticklabels([str(m) for m in ms])
    ax.set_xlabel("ordem m"); ax.set_ylabel("I/O médio por operação")
    ax.set_title("Resultados por operação — I/O por op (N=100k, aleatório)")
    ax.grid(True, axis="y", alpha=0.3); ax.legend(fontsize=11)
    return _save(fig, "chart_phase_io.png")


def _total_wall_by_m(rows):
    tot = defaultdict(float)
    for r in rows:
        if r["mode"] == "rand" and r["phase"] in ("insert", "search", "delete"):
            tot[r["m"]] += r["wall_ms"]
    return tot


def chart_wall_machines():
    """Comparação dos 2 sistemas: tempo de parede total vs ordem m."""
    nb, ti = _total_wall_by_m(ORDER), _total_wall_by_m(ORDER_LAP)
    ms = sorted(set(nb) & set(ti))
    fig, ax = plt.subplots(figsize=FIG, dpi=DPI)
    ax.plot(ms, [nb[m] for m in ms], "o-", color=ACCENT, lw=2.4, ms=6, label="Notebook")
    ax.plot(ms, [ti[m] for m in ms], "s-", color=TITAN, lw=2.4, ms=6, label="Titan (USP)")
    ax.set_xscale("log"); ax.set_xlabel("ordem m (escala log)")
    ax.set_ylabel("tempo total ins+busca+rem (ms)")
    ax.set_title("Dois sistemas — wall time vs m (N=100k, aleatório)")
    ax.grid(True, alpha=0.3); ax.legend(fontsize=11)
    return _save(fig, "chart_wall_machines.png")


def chart_cpu_split():
    """Comparação dos 2 sistemas: CPU de usuário vs CPU de sistema (I/O)."""
    import numpy as np
    ms = [3, 100, 1000]

    def split(rows):
        u = defaultdict(float); s = defaultdict(float)
        for r in rows:
            if r["mode"] == "rand" and r["phase"] in ("insert", "search", "delete"):
                u[r["m"]] += r["cpu_user_ms"]; s[r["m"]] += r["cpu_sys_ms"]
        return u, s

    nu, ns = split(ORDER); tu, ts = split(ORDER_LAP)
    x = np.arange(len(ms)); w = 0.36
    fig, ax = plt.subplots(figsize=FIG, dpi=DPI)
    ax.bar(x - w/2, [nu[m]/1000 for m in ms], w, color=ACCENT, label="Notebook · CPU usuário")
    ax.bar(x - w/2, [ns[m]/1000 for m in ms], w, bottom=[nu[m]/1000 for m in ms],
           color="#93c5fd", label="Notebook · CPU sistema (I/O)")
    ax.bar(x + w/2, [tu[m]/1000 for m in ms], w, color=TITAN, label="Titan · CPU usuário")
    ax.bar(x + w/2, [ts[m]/1000 for m in ms], w, bottom=[tu[m]/1000 for m in ms],
           color="#c4b5fd", label="Titan · CPU sistema (I/O)")
    ax.set_xticks(x); ax.set_xticklabels([f"m={m}" for m in ms])
    ax.set_ylabel("tempo de CPU total (s)")
    ax.set_title("Dois sistemas — CPU usuário vs sistema (N=100k)")
    ax.grid(True, axis="y", alpha=0.3); ax.legend(fontsize=8.5)
    return _save(fig, "chart_cpu_split.png")


def chart_determinism():
    """Prova de determinismo: acessos a disco idênticos nas 2 máquinas."""
    nb = {(r["m"], r["phase"]): r["disk_accesses"]
          for r in ORDER if r["mode"] == "rand"}
    ti = {(r["m"], r["phase"]): r["disk_accesses"]
          for r in ORDER_LAP if r["mode"] == "rand"}
    keys = sorted(set(nb) & set(ti))
    xs = [nb[k] for k in keys]; ys = [ti[k] for k in keys]
    fig, ax = plt.subplots(figsize=FIG, dpi=DPI)
    lim = max(xs + ys) * 1.05
    ax.plot([0, lim], [0, lim], "--", color=MUTED, lw=1.5, label="y = x (idêntico)")
    ax.scatter(xs, ys, s=55, color=ACCENT2, zorder=3, edgecolor="white")
    ax.set_xlim(0, lim); ax.set_ylim(0, lim)
    ax.set_xlabel("acessos a disco — Notebook")
    ax.set_ylabel("acessos a disco — Titan (USP)")
    ax.set_title("Determinismo: acessos a disco IDÊNTICOS nas 2 máquinas")
    ax.grid(True, alpha=0.3); ax.legend(fontsize=11)
    return _save(fig, "chart_determinism.png")


def chart_scale_machines():
    """Comparação dos 2 sistemas: escala do conjunto (wall vs N, m=100)."""
    def pts(rows):
        d = {}
        for r in rows:
            if r["m"] == 100 and r["mode"] == "rand" and r["phase"] == "search":
                d[r["N"]] = r["wall_ms"]
        return sorted(d.items())
    nb, ti = pts(SCALE), pts(SCALE_LAP)
    fig, ax = plt.subplots(figsize=FIG, dpi=DPI)
    if nb:
        ax.plot(*zip(*nb), "o-", color=ACCENT, lw=2.4, ms=6, label="Notebook")
    if ti:
        ax.plot(*zip(*ti), "s-", color=TITAN, lw=2.4, ms=6, label="Titan (USP)")
    ax.set_xscale("log"); ax.set_yscale("log")
    ax.set_xlabel("N — nº de chaves (escala log)")
    ax.set_ylabel("tempo de busca total (ms, escala log)")
    ax.set_title("Dois sistemas — escala do conjunto (m=100, busca)")
    ax.grid(True, alpha=0.3, which="both"); ax.legend(fontsize=11)
    return _save(fig, "chart_scale_machines.png")


def chart_node_anatomy():
    """Diagrama: como um nó da Árvore B é organizado (campos n, A[], K[])."""
    from matplotlib.patches import Rectangle
    fig, ax = plt.subplots(figsize=(9.0, 3.4), dpi=DPI)
    ax.set_xlim(0, 13); ax.set_ylim(0, 5); ax.axis("off")
    cells = [("n=3", INK, "#e0e7ff"),
             ("A0", ACCENT, "#dbeafe"), ("30", ACCENT2, "#fee2e2"),
             ("A1", ACCENT, "#dbeafe"), ("60", ACCENT2, "#fee2e2"),
             ("A2", ACCENT, "#dbeafe"), ("85", ACCENT2, "#fee2e2"),
             ("A3", ACCENT, "#dbeafe"), ("·", MUTED, "#f3f4f6"),
             ("A4", ACCENT, "#dbeafe")]
    x = 1.0; w = 1.1
    for txt, fg, bg in cells:
        ax.add_patch(Rectangle((x, 2.6), w, 1.2, facecolor=bg, edgecolor=INK, lw=1.3))
        ax.text(x + w/2, 3.2, txt, ha="center", va="center", color=fg,
                fontsize=12, fontweight="bold")
        x += w
    ax.text(1.0 + w/2, 4.1, "n", ha="center", color=INK, fontsize=10)
    ax.annotate("ponteiros A[ ] — RRN dos filhos no disco",
                xy=(3.2, 2.55), xytext=(3.2, 1.4), ha="center", color=ACCENT,
                fontsize=11, arrowprops=dict(arrowstyle="->", color=ACCENT))
    ax.annotate("chaves K[ ] — os valores armazenados",
                xy=(5.4, 3.85), xytext=(7.5, 4.6), ha="center", color=ACCENT2,
                fontsize=11, arrowprops=dict(arrowstyle="->", color=ACCENT2))
    ax.text(6.5, 0.5, "Registro de tamanho fixo: 1 inteiro (n) + m ponteiros + m chaves "
            "= PAGE_SIZE bytes  →  1 nó = 1 acesso ao disco",
            ha="center", color=MUTED, fontsize=10.5)
    ax.set_title("Anatomia de um nó (exemplo de ordem m=5: até 4 chaves, 5 filhos)",
                 fontsize=12)
    return _save(fig, "chart_node_anatomy.png")


def chart_file_layout():
    """Diagrama: o arquivo no disco como um vetor de registros de tamanho fixo."""
    from matplotlib.patches import Rectangle, FancyArrowPatch
    fig, ax = plt.subplots(figsize=(9.0, 3.4), dpi=DPI)
    ax.set_xlim(0, 13); ax.set_ylim(0, 5); ax.axis("off")
    recs = [("RRN 0\nHEADER\nroot·total·free_head", "#ede9fe", TITAN),
            ("RRN 1\nnó", "#dbeafe", ACCENT),
            ("RRN 2\nnó", "#dbeafe", ACCENT),
            ("RRN 3\nnó LIVRE", "#fee2e2", ACCENT2),
            ("RRN 4\nnó", "#dbeafe", ACCENT),
            ("…", "#f3f4f6", MUTED)]
    x = 0.6; w = 2.0
    centers = []
    for txt, bg, fg in recs:
        ax.add_patch(Rectangle((x, 2.2), w, 1.7, facecolor=bg, edgecolor=INK, lw=1.3))
        ax.text(x + w/2, 3.05, txt, ha="center", va="center", color=fg,
                fontsize=10.5, fontweight="bold")
        centers.append(x + w/2); x += w + 0.05
    ax.text(6.4, 1.4, "posição no arquivo:  offset(rrn) = rrn × PAGE_SIZE  "
            "(acesso direto, sem varrer nada)", ha="center", color=MUTED, fontsize=10.5)
    ax.add_patch(FancyArrowPatch((centers[0], 2.15), (centers[3], 2.15),
                 connectionstyle="arc3,rad=-0.35", arrowstyle="->",
                 color=ACCENT2, lw=1.6))
    ax.text((centers[0]+centers[3])/2, 0.7, "free_head → lista de nós livres "
            "(reaproveitamento)", ha="center", color=ACCENT2, fontsize=10)
    ax.set_title("Layout do arquivo: um vetor de registros de tamanho fixo",
                 fontsize=12)
    return _save(fig, "chart_file_layout.png")


print("gerando gráficos...")
C_NODE = chart_node_anatomy()
C_FILE = chart_file_layout()
C_PHASE = chart_phase_io()
C_WALL_M = chart_wall_machines()
C_CPU = chart_cpu_split()
C_DET = chart_determinism()
C_SCALE_M = chart_scale_machines()
C_IO_M = chart_io_vs_m()
C_H_M = chart_height_vs_m()
C_IO_N = chart_io_vs_N()
C_REUSE = chart_reuse()
C_MACH = chart_machines()

def order_table():
    by_m = defaultdict(dict)
    for r in ORDER:
        if r["mode"] == "rand":
            by_m[r["m"]][r["phase"]] = r
    headers = ["m", "altura", "busca I/O/op", "inser. I/O/op", "rem. I/O/op", "arquivo (KB)"]
    rows = []
    for m in sorted(by_m):
        d = by_m[m]
        rows.append([str(m), str(d["insert"]["height"]),
                     f'{d["search"]["avg_io_per_op"]:.2f}',
                     f'{d["insert"]["avg_io_per_op"]:.2f}',
                     f'{d["delete"]["avg_io_per_op"]:.2f}',
                     f'{d["insert"]["file_bytes"]/1024:,.0f}'])
    return headers, rows

def reuse_table():
    ref = [r for r in REUSE if r["phase"] == "churn_refill"]
    by = defaultdict(dict)
    for r in ref:
        by[(r["m"], r["N"])][r["reuse"]] = r
    headers = ["m", "N", "sem reuso (KB)", "com reuso (KB)", "economia"]
    rows = []
    for (m, N) in sorted(by):
        on = by[(m, N)].get("1"); off = by[(m, N)].get("0")
        if on and off and off["file_bytes"]:
            save = (1 - on["file_bytes"]/off["file_bytes"]) * 100
            rows.append([str(m), f"{N:,}", f'{off["file_bytes"]/1024:,.0f}',
                         f'{on["file_bytes"]/1024:,.0f}', f"{save:.0f}%"])
    return headers, rows

def machines_table():
    """Comparação lado a lado dos 2 sistemas (wall, CPU, acessos idênticos)."""
    def agg(rows):
        wall = defaultdict(float); cpu = defaultdict(float); acc = defaultdict(int)
        for r in rows:
            if r["mode"] == "rand" and r["phase"] in ("insert", "search", "delete"):
                wall[r["m"]] += r["wall_ms"]
                cpu[r["m"]] += r["cpu_user_ms"] + r["cpu_sys_ms"]
                acc[r["m"]] += r["disk_accesses"]
        return wall, cpu, acc
    nw, nc, na = agg(ORDER); tw, tc, ta = agg(ORDER_LAP)
    headers = ["m", "acessos disco", "Notebook wall", "Titan wall", "speedup"]
    rows = []
    for m in [3, 16, 100, 1000]:
        if m in nw and m in tw and tw[m]:
            speed = nw[m] / tw[m]
            ident = "✓ iguais" if na[m] == ta[m] else "DIVERGE"
            rows.append([str(m), f"{na[m]:,}".replace(",", "."),
                         f"{nw[m]/1000:.2f} s", f"{tw[m]/1000:.2f} s",
                         f"{speed:.2f}×"])
    return headers, rows


ORDER_HEADERS, ORDER_ROWS = order_table()
REUSE_HEADERS, REUSE_ROWS = reuse_table()
MACH_HEADERS, MACH_ROWS = machines_table()

def G(name):
    return os.path.join(GRAPHS, name)

SLIDES = [
    {
        "layout": "cover",
        "title": "Árvore B em Memória Secundária",
        "subtitle": "Implementação de uma classe Árvore B de ordem m operando estritamente em disco",
        "bullets": [
            "Algoritmos e Estruturas de Dados — AED-PG-2026 (5955001-3)",
            "Prof. Dr. José Augusto Baranauskas — DCM / USP",
            "Linguagem: C++17  ·  1º Semestre/2026",
            "Alunos: Cássio Takarada  ·  Cézio Luiz Ferreira Junior",
        ],
    },
    {
        "layout": "two_col",
        "title": "2. Roteiro da apresentação",
        "left_title": "Parte 1 — a estrutura",
        "left": [
            "Contexto: por que disco e por que Árvore B.",
            "Como um nó é organizado e como o arquivo é gravado.",
            "Decisões de implementação (1 nó por I/O).",
            "Os métodos: busca, inserção e remoção.",
            "Demonstração ao vivo + grafos gerados.",
        ],
        "right_title": "Parte 2 — a avaliação",
        "right": [
            "Experimentos, métricas e metodologia.",
            "Impacto da ordem m (I/O e altura).",
            "Resultados por operação e por tamanho do conjunto.",
            "Reaproveitamento de nós (espaço).",
            "Validação em dois sistemas + discussão (union) e conclusões.",
        ],
    },
    {
        "layout": "bullets",
        "title": "3. O problema e os objetivos",
        "bullets": [
            "Enunciado: implementar uma classe Árvore B de ordem m que resida em MEMÓRIA SECUNDÁRIA (um arquivo em disco).",
            "Restrição central: a árvore é lida/transferida em blocos — carregada parcialmente, NUNCA inteira na RAM. Cada operação toca um nó por vez.",
            "Operações exigidas: busca, inserção e remoção completas, mantendo a árvore balanceada.",
            "Instrumentação pedida: contar acessos ao disco e permitir monitorar a estrutura (impressão, altura, grafo).",
            "Nosso objetivo extra: avaliar empiricamente o efeito de m e do tamanho do conjunto, e validar os resultados em mais de uma máquina.",
        ],
    },
    {
        "layout": "bullets",
        "title": "4. Contexto — memória rápida × disco lento",
        "bullets": [
            "A memória RAM é rapidíssima, mas é pequena e volátil (some quando desliga).",
            "O disco é enorme e permanente, porém MILHARES de vezes mais lento que a RAM.",
            "Num banco de dados, o gargalo quase sempre é o disco — não a CPU.",
            "Por isso a regra do trabalho: a árvore mora no disco e nunca é carregada inteira; cada passo lê ou escreve só um pedaço (um nó).",
            "Consequência: para ser rápido, o segredo é fazer o MENOR número de acessos ao disco possível.",
        ],
    },
    {
        "layout": "bullets",
        "title": "5. Contexto — índice e a métrica honesta",
        "bullets": [
            "Um índice é o \"sumário\" que evita ler o arquivo inteiro: ele diz onde está cada registro.",
            "A Árvore B é a estrutura de índice usada de fato por bancos de dados (MySQL, PostgreSQL) e por sistemas de arquivos.",
            "A métrica que realmente importa NÃO é o relógio, e sim o NÚMERO DE ACESSOS AO DISCO — independe do hardware.",
            "Por que o relógio engana: a escrita vai para o cache do sistema operacional, então o tempo medido varia com a máquina e a carga.",
            "Toda a apresentação gira em torno de uma pergunta: como organizar o arquivo para responder buscas com o menor número de acessos?",
        ],
    },
    {
        "layout": "text_image",
        "title": "6. Anatomia de um nó",
        "bullets": [
            "Cada nó é um registro de TAMANHO FIXO, gravado e lido de uma vez só.",
            "Campos: n (quantas chaves o nó tem), o vetor de ponteiros A[ ] (os RRN dos filhos) e o vetor de chaves K[ ] (os valores).",
            "Um nó de ordem m guarda até m-1 chaves e até m filhos — por isso m maior = mais chaves por nó.",
            "PAGE_SIZE = 1 inteiro + m ponteiros + m chaves. Esse tamanho fixo é o que permite achar qualquer nó por cálculo direto.",
            "Regra de ouro: ler ou escrever UM nó = UM acesso ao disco.",
        ],
        "images": [(C_NODE, None)],
    },
    {
        "layout": "text_image",
        "title": "7. Layout do arquivo no disco",
        "bullets": [
            "O arquivo é um VETOR de registros de tamanho fixo, numerados por RRN (Relative Record Number).",
            "O registro 0 é o HEADER: guarda a raiz (root), o total de nós e a cabeça da lista de livres (free_head).",
            "Para achar um nó: offset = RRN × PAGE_SIZE. Acesso direto, sem varrer o arquivo.",
            "Nós removidos não somem: entram numa lista de livres encadeada no próprio arquivo, prontos para reuso.",
            "Abrir o arquivo e já saber a raiz (no header) custa uma única leitura de metadado.",
        ],
        "images": [(C_FILE, None)],
    },
    {
        "layout": "image_full",
        "title": "8. Decisões de implementação",
        "bullets": [
            "Ordem m é constante simbólica de compilação (M): estrutura totalmente parametrizada.",
            "Raiz da árvore: armazenada no header (registro 0) do arquivo — campo root.",
            "Um nó por I/O: readNode/writeNode de um registro de PAGE_SIZE; a árvore NUNCA é carregada inteira em memória.",
            "Registro do nó: n  ·  A[0..m-1] (ponteiros RRN)  ·  K[0..m-1] (chaves).",
            "Reaproveitamento de nós: estrutura = lista encadeada de livres (pilha LIFO) no próprio arquivo - free_head no header; cada nó livre marca n=-1 e aponta o próximo em K[1]; allocNode reusa antes de estender o arquivo.",
        ],
        "images": [(G("m3_final.png"), "Árvore B (m=3) com 40 chaves — exportada do nosso código (Graphviz); cada caixa mostra o RRN do nó em disco")],
    },
    {
        "layout": "bullets",
        "title": "9. O DiskManager — um nó por I/O",
        "bullets": [
            "Toda a conversa com o disco passa por uma única classe: o DiskManager.",
            "readNode(rrn) / writeNode(rrn): leem e escrevem exatamente UM registro de PAGE_SIZE — nunca mais que isso.",
            "É aqui que mora o CONTADOR DE ACESSOS: cada readNode/writeNode soma 1. Essa é a métrica central da avaliação.",
            "readHeader / writeHeader: acessam o registro 0 (raiz, total, free_head).",
            "Garantia de projeto: a árvore nunca é carregada inteira — a lógica da B-tree só enxerga o disco por meio do DiskManager.",
        ],
    },
    {
        "layout": "bullets",
        "title": "10. Reaproveitamento — a lista de livres",
        "bullets": [
            "Quando uma remoção libera um nó, ele não é descartado: vira um nó LIVRE.",
            "freeNode(rrn): marca o nó com n = -1 e o encadeia na lista de livres, cuja cabeça (free_head) fica no header.",
            "allocNode(): antes de crescer o arquivo, pega o primeiro nó da lista de livres (pilha LIFO) — só estende o arquivo se não houver livres.",
            "Resultado: o arquivo reaproveita espaço em vez de inchar a cada remoção+inserção.",
            "É um comportamento ligável (setReuse) — usamos isso para medir o ganho no experimento de churn.",
        ],
    },
    {
        "layout": "bullets",
        "title": "11. Métodos — visão geral",
        "bullets": [
            "Três operações principais, todas operando em disco: busca, inserção e remoção.",
            "mSearch / mSearchPath — busca m-way da raiz à folha; retorna (RRN, índice, achou).",
            "insertB — inserção bottom-up com propagação de split.",
            "deleteB — remoção com sucessor in-order, redistribuição (empréstimo) e fusão (merge).",
            "Auxiliares: splitNode, insertInNode, removeFromNode, findSuccessorLeaf.",
            "Monitoramento: printTree (hierárquico), height, exportDot (Graphviz) e o contador de acessos.",
            "Nos próximos slides, cada operação em detalhe.",
        ],
    },
    {
        "layout": "text_image",
        "title": "12. Busca m-way",
        "bullets": [
            "Em cada nó, as chaves dividem o universo em FAIXAS — como um dicionário que diz 'entre tal e tal palavra'.",
            "Comparamos a chave buscada com as chaves do nó e descemos pelo ponteiro da faixa certa.",
            "Repetimos nó a nó até achar a chave ou chegar a uma folha sem ela.",
            "Custo: aproximadamente a ALTURA da árvore em acessos ao disco — pouquíssimos, porque a árvore é rasa.",
            "No grafo: buscar 70 desce da raiz (40, 60) para o filho da direita — 2 acessos.",
        ],
        "images": [(G("m3_final.png"), "Árvore B de ordem 3 — a busca segue um caminho da raiz à folha")],
    },
    {
        "layout": "bullets",
        "title": "13. Inserção e split",
        "bullets": [
            "A inserção é bottom-up: primeiro a busca encontra a FOLHA correta, depois inserimos a chave em ordem.",
            "Se o nó couber, acabou. Se ESTOURAR (passa de m-1 chaves), fazemos o split.",
            "Split: o nó cheio é dividido em dois e a chave do MEIO (mediana) SOBE para o nó pai.",
            "Se o pai também estourar, o split se propaga para cima — e pode até criar uma nova raiz, aumentando a altura.",
            "É exatamente esse mecanismo que mantém a árvore sempre BALANCEADA, sem precisar reorganizar tudo.",
        ],
    },
    {
        "layout": "bullets",
        "title": "14. Remoção — sucessor, redistribuição e fusão",
        "bullets": [
            "É a operação mais delicada. Se a chave está num nó INTERNO, trocamos pelo seu sucessor in-order (a menor chave à direita), reduzindo o caso ao de uma folha.",
            "Ao remover de uma folha, o nó pode ficar ABAIXO DO MÍNIMO de chaves (underflow).",
            "Reparo 1 — redistribuição: se um irmão tem chaves sobrando, pegamos uma emprestada (passando pelo pai).",
            "Reparo 2 — fusão: se nenhum irmão pode emprestar, fundimos o nó com um irmão e puxamos uma chave do pai.",
            "A fusão pode propagar o underflow para cima — no limite, a raiz encolhe e a árvore fica mais baixa.",
        ],
    },
    {
        "layout": "two_col",
        "title": "15. Resumo — tudo que o código faz",
        "left_title": "Núcleo — Árvore B em disco",
        "left": [
            "Árvore B de ordem m parametrizável (M em tempo de compilação).",
            "Opera 100% em disco: 1 nó por I/O; nada é carregado em RAM.",
            "Busca m-way (mSearch / mSearchPath).",
            "Inserção bottom-up com split (insertB).",
            "Remoção com sucessor, redistribuição e fusão (deleteB).",
            "Persistência: header (root, total, free_head) + nós de tamanho fixo, acesso direto por RRN.",
        ],
        "right_title": "Extras / funcionalidades acessórias",
        "right": [
            "Contador de acessos ao disco (métrica central da avaliação).",
            "Reaproveitamento de nós via free list, com liga/desliga (setReuse).",
            "Impressão hierárquica (printTree) e cálculo de altura (height).",
            "Exportação Graphviz (exportDot) -> imagens PNG dos grafos.",
            "Medição de ocupação: total_nodes, free_nodes, fileSizeBytes.",
            "Benchmark (bench): CSV com I/O, tempo CPU/IO (getrusage), altura, nós, tamanho.",
            "Automação: run_all.sh + make_tables.py (tabelas/gráficos) + compare.py (2 máquinas).",
            "Testes de stress + menu interativo (CRUD + métricas).",
        ],
    },
    {
        "layout": "demo",
        "title": "16. Demonstração — o programa rodando",
        "bullets": [
            "Animação real: do `make M=3` (compilação) até cada item do menu da main.",
            "Inserir · Buscar · Remover · ver Acessos ao disco · Imprimir árvore · Exportar grafo.",
            "Repetimos com m=5: nós comportam mais chaves, a árvore fica mais rasa.",
            "Tudo capturado da execução de verdade — sem montagem.",
        ],
        "gif": os.path.join(ASSETS, "demo_btree.gif"),
        "poster": os.path.join(ASSETS, "demo_poster.png"),
        "caption": "make M=3/M=5 → menu interativo → grafo exportado pelo próprio programa",
    },
    {
        "layout": "image_full",
        "title": "17. Grafos gerados pelo programa",
        "bullets": [
            "O programa exporta a árvore para Graphviz (exportDot) e gera a imagem do grafo.",
            "Cada caixa é um nó GRAVADO NO DISCO — o número é o RRN (a posição do registro no arquivo).",
            "m=3: no máximo 2 chaves por nó → mais níveis. m=5: até 4 chaves por nó → árvore mais rasa, menos acessos.",
        ],
        "images": [(os.path.join(ASSETS, "graphs_m3_m5.png"),
                    "Mesma sequência de chaves, duas ordens m — à direita a árvore é mais rasa")],
    },
    {
        "layout": "text_image",
        "title": "18. Os experimentos — a matriz de testes",
        "bullets": [
            "Exp. 1 - Impacto da ordem m: m em {3,4,5,8,16,32,64,100,128,256,512,1000}, N=10^5.",
            "Exp. 2 - Escala do conjunto: N em {10^3,10^4,10^5,10^6}, m em {3,100,1000}.",
            "Exp. 3 - Ocupação do arquivo: churn COM e SEM reaproveitamento.",
            "Exp. 4 - Tempo: CPU (user+sys) vs espera de I/O (getrusage).",
            "Modos aleatório e sequencial em todos.",
            "Novidade: validação em 2 máquinas (Notebook x Titan/USP) - disco idêntico.",
        ],
        "images": [(C_IO_N, None)],
    },
    {
        "layout": "bullets",
        "title": "19. Metodologia e ambiente",
        "bullets": [
            "Driver não interativo (bench): roda inserção → busca → remoção e emite um CSV com todas as métricas.",
            "Cada configuração é reconstruída do zero (arquivo novo), garantindo medições independentes.",
            "Dois modos de carga: ALEATÓRIO (chaves embaralhadas) e SEQUENCIAL (em ordem) — estressam a árvore de formas diferentes.",
            "Tempo medido com getrusage: separa CPU de usuário, CPU de sistema (syscalls de I/O) e espera de I/O.",
            "Duas máquinas: meu Notebook e o servidor Titan (USP). Automação por run_all.sh + scripts Python.",
            "Métrica de referência sempre o contador de acessos — determinístico e comparável entre máquinas.",
        ],
    },
    {
        "layout": "text_image",
        "title": "20. Métricas utilizadas",
        "bullets": [
            "Acessos ao disco - contador readNode+writeNode; total e média/op (métrica honesta de I/O).",
            "Altura da árvore - nº de níveis (~ log_m N).",
            "Ocupação - nós físicos, nós livres e tamanho em bytes.",
            "Tempo de parede (wall) - relógio de alta resolução.",
            "Tempo de CPU - usuário (lógica) e sistema (syscalls de I/O), via getrusage.",
            "Espera de I/O - wall menos CPU.",
        ],
        "images": [(C_H_M, None)],
    },
    {
        "layout": "table_image",
        "title": "21. Impacto da ordem m — I/O por busca (N=10^5)",
        "table": (ORDER_HEADERS, ORDER_ROWS),
        "images": [(C_IO_M, None)],
        "footnote": "I/O por busca cai de 13,25 (m=3, altura 14) para 2,00 (m>=512, altura 2). Ganho satura por volta de m~64-128.",
    },
    {
        "layout": "text_image",
        "title": "22. Impacto da ordem m — altura da árvore",
        "bullets": [
            "A altura é o número de níveis — e quase iguala o número de acessos por busca.",
            "Com m=3, a árvore tem 14 níveis para 100 mil chaves. É muito 'fundo' para uma estrutura de disco.",
            "Aumentando m, a altura DESPENCA em degraus: poucos níveis bastam, porque cada nó comporta muito mais chaves.",
            "A partir de m grande, a altura chega a 2 — não dá para ficar mais raso.",
            "Altura baixa = poucos acessos = a razão de existir da Árvore B.",
        ],
        "images": [(C_H_M, None)],
    },
    {
        "layout": "text_image",
        "title": "23. Resultados separados por operação",
        "bullets": [
            "Mesmo experimento, agora separando inserção, busca e remoção.",
            "A BUSCA é a mais barata: desce direto da raiz à folha (~altura acessos).",
            "INSERÇÃO custa um pouco mais: além de descer, pode reescrever nós no caminho (split).",
            "REMOÇÃO é a mais cara: pode ler irmãos para redistribuir/fundir nós (reparo de underflow).",
            "As três caem junto com m — e todas saturam quando a árvore chega a 2-3 níveis.",
            "Ou seja: o ganho de aumentar m vale para TODAS as operações, não só para a busca.",
        ],
        "images": [(C_PHASE, None)],
    },
    {
        "layout": "text_image",
        "title": "24. Escala do conjunto (N de mil a 1 milhão)",
        "bullets": [
            "Agora fixamos a ordem e variamos o tamanho do conjunto N, de mil até um MILHÃO de chaves.",
            "Mesmo multiplicando N por mil, o número de acessos por busca quase não muda.",
            "Motivo: o custo é logarítmico na base m — dobrar/multiplicar os dados acrescenta no máximo um nível.",
            "Com m grande, a curva fica praticamente plana: a estrutura aguenta bases enormes.",
            "É a tradução prática de 'escalável': cresce o dado, não cresce o custo por operação.",
        ],
        "images": [(C_IO_N, None)],
    },
    {
        "layout": "table_image",
        "title": "25. Reaproveitamento de nós — resultado (churn)",
        "table": (REUSE_HEADERS, REUSE_ROWS),
        "images": [(C_REUSE, None)],
        "footnote": "O reaproveitamento de nós economiza ~27-30% do tamanho do arquivo neste workload.",
    },
    {
        "layout": "table_image",
        "title": "26. Dois sistemas — comparação de tempo (wall)",
        "table": (MACH_HEADERS, MACH_ROWS),
        "images": [(C_WALL_M, None)],
        "footnote": "Mesmo programa, mesmo arquivo, duas máquinas. Acessos a disco idênticos; só o tempo muda — o Titan (USP) é mais rápido, mas a CURVA tem a mesma forma.",
    },
    {
        "layout": "two_image",
        "title": "27. Dois sistemas — CPU e prova de determinismo",
        "bullets": [
            "Esquerda: o tempo é quase todo CPU de SISTEMA (chamadas de I/O ao SO), não espera de disco — efeito do cache de páginas. Por isso o relógio não é a métrica honesta.",
            "Direita: cada ponto é (acessos no Notebook, acessos no Titan). TODOS caem na reta y=x — ou seja, o número de acessos é exatamente igual nas duas máquinas.",
            "Conclusão: a implementação é determinística e portável; o hardware muda o tempo, nunca o comportamento lógico.",
        ],
        "images": [(C_CPU, "CPU usuário vs sistema (I/O) nas 2 máquinas"),
                   (C_DET, "Acessos a disco idênticos — pontos sobre a reta y=x")],
    },
    {
        "layout": "text_image",
        "title": "28. Dois sistemas — validação cruzada (resumo)",
        "bullets": [
            "Execução em duas máquinas: Notebook x Titan (USP).",
            "Acessos a disco IDÊNTICOS nas duas — resultado determinístico, independente do hardware.",
            "Apenas o tempo varia: wall time (inclui espera de I/O e carga) e tempo de CPU (user+sys, só o trabalho do processo).",
            "Comparação automatizada via compare.py a partir dos CSVs de cada máquina.",
        ],
        "images": [(C_MACH, None)] if C_MACH else [],
    },
    {
        "layout": "text_image",
        "title": "29. Dois sistemas — escala do conjunto",
        "bullets": [
            "Agora variando N de mil a um MILHÃO de chaves (ordem m=100), nas duas máquinas.",
            "Multiplicamos N por 1000 e o tempo de busca cresce devagar — porque a árvore é logarítmica na base m.",
            "As duas máquinas mantêm a mesma tendência; o Titan fica abaixo (mais rápido), em paralelo.",
            "É a prova prática de escalabilidade: a estrutura aguenta bases grandes sem explodir o custo.",
        ],
        "images": [(C_SCALE_M, None)],
    },
    {
        "layout": "two_col",
        "title": "30. E se usássemos union no header? (sugestão do prof.)",
        "left_title": "Como está hoje",
        "left": [
            "Header { root, total, free_head } e BNode { n, A[], K[] } são dois tipos SEPARADOS.",
            "Toda E/S serializa campo a campo com memcpy (readHeader/writeHeader/readNode/writeNode).",
            "O registro 0 (header) usa só 12 bytes, mas reserva uma página inteira (PAGE_SIZE) para manter o RRN uniforme: offset = rrn × PAGE_SIZE.",
            "O nó LIVRE \"esconde\" o ponteiro do próximo em K[1] (com n = -1) — funciona, mas é implícito.",
        ],
        "right_title": "Com union { Header; BNode; FreeNode; }",
        "right": [
            "Um único tipo Página: todos os registros (header e nós) compartilham a MESMA página de tamanho fixo — garantido em tempo de compilação por sizeof.",
            "Lê/escreve a página inteira de uma vez e a interpreta como header (RRN 0) ou nó (RRN ≥ 1): some a serialização campo a campo → menos código e menos bugs de offset (justo as dificuldades que tivemos).",
            "A lista de livres ganha um campo explícito FreeNode.next no lugar do \"K[1] mágico\" → autodocumentado.",
            "Custo: gravar struct cru depende de layout/padding/endianness → menos portável (a validação Notebook×Titan supõe layout igual); ler membro ≠ do escrito é UB em C++ — mitiga-se com memcpy / std::bit_cast.",
            "Veredito: página uniforme e código enxuto, em troca de cuidado com a portabilidade do binário.",
        ],
    },
    {
        "layout": "bullets",
        "title": "31. Dificuldades técnicas",
        "bullets": [
            "Indexação 1-based dos nós (RRN 0 é o header): exigiu cuidado em todo cálculo de posição.",
            "Remoção: manter o caminho (path) na ordem certa para reparar o underflow foi a parte mais difícil.",
            "eofbit do fstream: uma leitura no fim do arquivo deixava o stream em estado de erro e fazia leituras seguintes falharem em silêncio — resolvido com clear().",
            "Garantir de verdade 1 nó por I/O: nenhuma estrutura podia 'esconder' a árvore em RAM.",
            "Serialização campo a campo (memcpy): funciona, mas é verbosa e propensa a erros de offset — daí a sugestão do union.",
        ],
    },
    {
        "layout": "two_col",
        "title": "32. Vantagens × Desvantagens (trade-offs)",
        "left_title": "Vantagens",
        "left": [
            "Altura baixa → pouquíssimos acessos ao disco.",
            "Sempre balanceada, sem reorganização global.",
            "Ideal para memória secundária e índices de SGBD.",
            "Determinística: mesmo resultado em qualquer máquina.",
            "Reaproveitamento de espaço via lista de livres.",
        ],
        "right_title": "Desvantagens / custos",
        "right": [
            "Nós podem ficar subocupados (~50% no caso sequencial).",
            "Remoção é complexa (sucessor, redistribuição, fusão).",
            "Para m muito grande, varrer as chaves dentro do nó passa a custar CPU.",
            "Trade-off central: m maior → menos I/O, mas mais CPU por nó.",
            "Existe, portanto, um m ÓTIMO (nó do tamanho de um bloco de disco).",
        ],
    },
    {
        "layout": "bullets",
        "title": "33. Aplicações práticas",
        "bullets": [
            "Índices de SGBDs: B-tree / B+ tree são a base dos índices de MySQL/InnoDB, PostgreSQL, Oracle e SQL Server.",
            "Sistemas de arquivos: NTFS, HFS+/APFS, ext4 (HTree) e Btrfs usam B-trees para diretórios e metadados.",
            "Armazenamento chave-valor / embarcado: BerkeleyDB, SQLite, LMDB.",
            "Onde eu usaria: como índice (primário ou secundário) de um banco de dados em disco - exatamente o cenário deste trabalho.",
        ],
    },
    {
        "layout": "two_col",
        "title": "34. Ferramentas utilizadas",
        "left_title": "Implementação & experimentos",
        "left": [
            "Linguagem: C++17 (g++ -O2) — núcleo da Árvore B em disco.",
            "Build: GNU Make (parametriza a ordem via -DM=<m>).",
            "Persistência: std::fstream (arquivo binário, registros de tamanho fixo).",
            "Graphviz (dot): renderiza os grafos exportados pelo programa (exportDot).",
            "Python 3 + matplotlib / numpy: gráficos dos experimentos.",
            "Pillow (PIL): montagem dos GIFs da demonstração.",
            "python-pptx + reportlab: geração automática dos slides (PPTX e PDF).",
            "Bash: run_all.sh automatiza a suíte; getrusage mede CPU/I/O.",
            "Git/GitHub: versionamento. Validação cruzada: Notebook + servidor Titan (USP).",
        ],
        "right_title": "Processo & reprodutibilidade",
        "right": [
            "Versionamento com Git/GitHub ao longo de todo o desenvolvimento.",
            "Experimentos reprodutíveis: driver não interativo gera CSV de métricas.",
            "Build parametrizável pela ordem: make M=<m> reconstrói tudo do zero.",
            "Validação cruzada em duas máquinas (Notebook + servidor Titan/USP).",
            "Grafos exportados pelo próprio programa (exportDot → Graphviz).",
        ],
    },
    {
        "layout": "bullets",
        "title": "35. Conclusão",
        "bullets": [
            "Implementamos uma Árvore B que opera estritamente em disco, com 1 nó por I/O — exatamente como o enunciado pede.",
            "Achado 1: o I/O por operação cai com m até SATURAR (~ m=64-128). Existe um m ótimo — o nó dimensionado para um bloco de disco.",
            "Achado 2: o reaproveitamento de nós (free list) economiza ~27-30% do tamanho do arquivo no workload de churn.",
            "Achado 3: resultados DETERMINÍSTICOS — acessos a disco idênticos em 2 máquinas (Notebook x Titan/USP); só o tempo varia.",
            "Mensagem final: a Árvore B é rasa, balanceada e econômica em acessos — por isso é a base de bancos de dados e sistemas de arquivos.",
        ],
    },
    {
        "layout": "bullets",
        "title": "36. Referências",
        "bullets": [
            "Comer, D. (1979). The Ubiquitous B-Tree. ACM Computing Surveys, 11(2), 121-137.",
            "Knuth, D. The Art of Computer Programming, Vol. 3 — Sorting and Searching.",
            "Folk, M. & Zoellick, B. File Structures.",
            "Bayer, R. & McCreight, E. (1972). Organization and Maintenance of Large Ordered Indexes.",
            "Baranauskas, J. A. Slides AED-PG-2026 (Árvores B). DCM / USP.",
            "Obrigado! Ficamos à disposição para perguntas.",
        ],
    },
]

def build_pptx(path):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from PIL import Image

    def rgb(h):
        return RGBColor.from_string(h.lstrip("#"))

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    SW = prs.slide_width
    blank = prs.slide_layouts[6]

    def add_text(slide, left, top, width, height, runs, anchor=MSO_ANCHOR.TOP, space=6):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        first = True
        for (text, size, color, bold, bullet) in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False; p.space_after = Pt(space)
            r = p.add_run(); r.text = ("•  " + text) if bullet else text
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = rgb(color); r.font.name = "Calibri"
        return tb

    def add_table(slide, headers, rows, left, top, width, height):
        gtbl = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
        for j, h in enumerate(headers):
            c = gtbl.cell(0, j); c.text = h
            pr = c.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
            run = pr.runs[0]; run.font.size = Pt(11); run.font.bold = True
            run.font.color.rgb = rgb("#ffffff")
            c.fill.solid(); c.fill.fore_color.rgb = rgb(ACCENT)
        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                c = gtbl.cell(i, j); c.text = str(val)
                pr = c.text_frame.paragraphs[0]
                pr.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                run = pr.runs[0]; run.font.size = Pt(10); run.font.color.rgb = rgb(INK)
                c.fill.solid(); c.fill.fore_color.rgb = rgb("#eef2ff" if i % 2 == 0 else "#ffffff")
        return gtbl

    def add_image_fit(slide, img, left, top, max_w, max_h, caption=None):
        try:
            iw, ih = Image.open(img).size
        except Exception:
            iw, ih = (1280, 720)
        scale = min(max_w / iw, max_h / ih)
        w, h = int(iw * scale), int(ih * scale)
        l = left + (max_w - w) // 2
        slide.shapes.add_picture(img, l, top, width=w, height=h)
        if caption:
            add_text(slide, left, top + h + Emu(30000), max_w, Inches(0.4),
                     [(caption, 9, MUTED, False, False)])

    def title_band(slide, text):
        band = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.95))
        band.fill.solid(); band.fill.fore_color.rgb = rgb(ACCENT); band.line.fill.background()
        add_text(slide, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.72),
                 [(text, 23, "#ffffff", True, False)], anchor=MSO_ANCHOR.MIDDLE)

    def footer(slide, idx):
        add_text(slide, Inches(0.5), Inches(7.06), Inches(12.3), Inches(0.34),
                 [(f"Árvore B em Memória Secundária — AED-PG-2026          {idx+1}/{len(SLIDES)}",
                   9, MUTED, False, False)])

    for idx, s in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank)
        lay = s["layout"]

        if lay == "cover":
            band = slide.shapes.add_shape(1, 0, Inches(2.2), SW, Inches(1.7))
            band.fill.solid(); band.fill.fore_color.rgb = rgb(ACCENT); band.line.fill.background()
            add_text(slide, Inches(0.8), Inches(2.35), Inches(11.7), Inches(1.4),
                     [(s["title"], 40, "#ffffff", True, False)], anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.9),
                     [(s["subtitle"], 18, INK, False, False)])
            add_text(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.0),
                     [(b, 15, MUTED, False, False) for b in s["bullets"]])
            continue

        title_band(slide, s["title"])
        body_top = Inches(1.15)

        if lay == "image_full":
            add_text(slide, Inches(0.5), body_top, Inches(12.3), Inches(2.95),
                     [(b, 15, INK, False, True) for b in s["bullets"]], space=8)
            img, cap = s["images"][0]
            add_image_fit(slide, img, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.65), cap)

        elif lay == "two_col":
            lruns = [(s["left_title"], 15, ACCENT, True, False)] + \
                    [(b, 14, INK, False, True) for b in s["left"]]
            rruns = [(s["right_title"], 15, ACCENT, True, False)] + \
                    [(b, 14, INK, False, True) for b in s["right"]]
            add_text(slide, Inches(0.45), body_top, Inches(6.15), Inches(5.9), lruns, space=5)
            add_text(slide, Inches(6.85), body_top, Inches(6.15), Inches(5.9), rruns, space=5)

        elif lay == "text_image":
            add_text(slide, Inches(0.5), body_top, Inches(6.4), Inches(5.8),
                     [(b, 16, INK, False, True) for b in s["bullets"]],
                     anchor=MSO_ANCHOR.MIDDLE, space=8)
            img, cap = s["images"][0]
            add_image_fit(slide, img, Inches(6.95), Inches(1.5), Inches(6.1), Inches(5.0), cap)

        elif lay == "table_image":
            headers, rows = s["table"]
            add_table(slide, headers, rows, Inches(0.4), body_top, Inches(5.9), Inches(5.0))
            img, cap = s["images"][0]
            add_image_fit(slide, img, Inches(6.55), Inches(1.4), Inches(6.5), Inches(5.1), cap)
            if s.get("footnote"):
                add_text(slide, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.5),
                         [(s["footnote"], 11, MUTED, False, False)])

        elif lay == "demo":
            add_text(slide, Inches(0.5), body_top, Inches(12.3), Inches(1.65),
                     [(b, 15, INK, False, True) for b in s["bullets"]], space=6)
            add_image_fit(slide, s["gif"], Inches(2.5), Inches(3.0),
                          Inches(8.3), Inches(3.7), s.get("caption"))

        elif lay == "two_image":
            add_text(slide, Inches(0.5), body_top, Inches(12.3), Inches(2.05),
                     [(b, 14, INK, False, True) for b in s["bullets"]], space=4)
            (img1, cap1), (img2, cap2) = s["images"][0], s["images"][1]
            add_image_fit(slide, img1, Inches(0.4), Inches(3.5), Inches(6.2), Inches(3.2), cap1)
            add_image_fit(slide, img2, Inches(6.85), Inches(3.5), Inches(6.2), Inches(3.2), cap2)

        else:
            add_text(slide, Inches(0.7), body_top, Inches(12.0), Inches(5.8),
                     [(b, 20, INK, False, True) for b in s["bullets"]],
                     anchor=MSO_ANCHOR.MIDDLE, space=12)

        footer(slide, idx)

    prs.save(path)
    print("PPTX:", path)

def build_pdf(path):
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    from reportlab.lib.colors import HexColor, white

    W, H = 960, 540
    c = canvas.Canvas(path, pagesize=(W, H))

    def wrap(text, font, size, max_w):
        words = text.split(" "); lines, cur = [], ""
        for w in words:
            t = (cur + " " + w).strip()
            if c.stringWidth(t, font, size) <= max_w:
                cur = t
            else:
                if cur:
                    lines.append(cur)
                cur = w
        if cur:
            lines.append(cur)
        return lines

    def bullets(items, x, y, max_w, size=14, lh=20, gap=6):
        for b in items:
            lines = wrap(b, "Helvetica", size, max_w)
            c.setFillColor(HexColor(ACCENT)); c.setFont("Helvetica", size)
            c.drawString(x, y, "•")
            c.setFillColor(HexColor(INK))
            for ln in lines:
                c.drawString(x + 16, y, ln); y -= lh
            y -= gap
        return y

    def draw_image_fit(img, x, y_top, max_w, max_h, caption=None):
        try:
            ir = ImageReader(img); iw, ih = ir.getSize()
        except Exception:
            return
        scale = min(max_w / iw, max_h / ih)
        w, h = iw * scale, ih * scale
        c.drawImage(ir, x + (max_w - w) / 2, y_top - h, width=w, height=h,
                    preserveAspectRatio=True, mask="auto")
        if caption:
            c.setFont("Helvetica", 8.5); c.setFillColor(HexColor(MUTED))
            for k, ln in enumerate(wrap(caption, "Helvetica", 8.5, max_w)):
                c.drawCentredString(x + max_w / 2, y_top - h - 12 - k * 11, ln)

    def draw_table(headers, rows, x, y_top, total_w):
        nc = len(headers); cw = total_w / nc; rh = 21
        c.setFillColor(HexColor(ACCENT)); c.rect(x, y_top - rh, total_w, rh, fill=1, stroke=0)
        c.setFillColor(white); c.setFont("Helvetica-Bold", 9)
        for j, hd in enumerate(headers):
            c.drawCentredString(x + cw * j + cw / 2, y_top - rh + 7, str(hd))
        c.setFont("Helvetica", 9); yy = y_top - rh
        for i, row in enumerate(rows):
            if i % 2 == 0:
                c.setFillColor(HexColor("#eef2ff")); c.rect(x, yy - rh, total_w, rh, fill=1, stroke=0)
            c.setFillColor(HexColor(INK))
            for j, val in enumerate(row):
                if j == 0:
                    c.drawString(x + 6, yy - rh + 7, str(val))
                else:
                    c.drawCentredString(x + cw * j + cw / 2, yy - rh + 7, str(val))
            yy -= rh
        return yy

    for idx, s in enumerate(SLIDES):
        lay = s["layout"]
        c.setFillColor(HexColor("#ffffff")); c.rect(0, 0, W, H, fill=1, stroke=0)

        if lay == "cover":
            c.setFillColor(HexColor(ACCENT)); c.rect(0, H - 250, W, 130, fill=1, stroke=0)
            c.setFillColor(white); c.setFont("Helvetica-Bold", 34)
            for k, ln in enumerate(wrap(s["title"], "Helvetica-Bold", 34, W - 120)):
                c.drawString(60, H - 175 - k * 38, ln)
            c.setFillColor(HexColor(INK)); c.setFont("Helvetica", 15)
            for k, ln in enumerate(wrap(s["subtitle"], "Helvetica", 15, W - 120)):
                c.drawString(60, H - 285 - k * 20, ln)
            c.setFillColor(HexColor(MUTED)); c.setFont("Helvetica", 13)
            for k, b in enumerate(s["bullets"]):
                c.drawString(60, H - 360 - k * 24, b)
            c.showPage(); continue

        c.setFillColor(HexColor(ACCENT)); c.rect(0, H - 60, W, 60, fill=1, stroke=0)
        c.setFillColor(white); c.setFont("Helvetica-Bold", 18)
        c.drawString(36, H - 40, s["title"])
        top = H - 88

        if lay == "image_full":
            bullets(s["bullets"], 36, top, W - 90, size=14, lh=20, gap=7)
            img, cap = s["images"][0]
            draw_image_fit(img, 30, 220, W - 60, 180, cap)

        elif lay == "two_col":
            c.setFillColor(HexColor(ACCENT)); c.setFont("Helvetica-Bold", 13)
            c.drawString(36, top, s["left_title"]); c.drawString(W/2 + 10, top, s["right_title"])
            bullets(s["left"], 36, top - 24, W/2 - 70, size=11.5, lh=16, gap=7)
            bullets(s["right"], W/2 + 10, top - 24, W/2 - 60, size=11.5, lh=16, gap=7)

        elif lay == "text_image":
            bullets(s["bullets"], 36, top - 30, 442, size=14, lh=22, gap=9)
            img, cap = s["images"][0]
            draw_image_fit(img, 500, top + 6, 430, 380, cap)

        elif lay == "table_image":
            headers, rows = s["table"]
            draw_table(headers, rows, 28, top, 430)
            img, cap = s["images"][0]
            draw_image_fit(img, 470, top, 470, 380, cap)
            if s.get("footnote"):
                c.setFillColor(HexColor(MUTED)); c.setFont("Helvetica-Oblique", 10)
                for k, ln in enumerate(wrap(s["footnote"], "Helvetica-Oblique", 10, W - 70)):
                    c.drawString(36, 56 - k * 13, ln)

        elif lay == "demo":
            bullets(s["bullets"], 36, top, W - 90, size=12, lh=16, gap=4)
            draw_image_fit(s.get("poster", s.get("gif")), 30, 300, W - 60, 250,
                           s.get("caption"))

        elif lay == "two_image":
            bullets(s["bullets"], 36, top, W - 90, size=11, lh=14.5, gap=4)
            (img1, cap1), (img2, cap2) = s["images"][0], s["images"][1]
            draw_image_fit(img1, 18, 250, W / 2 - 28, 200, cap1)
            draw_image_fit(img2, W / 2 + 10, 250, W / 2 - 28, 200, cap2)

        else:
            bullets(s["bullets"], 40, top - 60, W - 110, size=17, lh=27, gap=10)

        c.setFillColor(HexColor(MUTED)); c.setFont("Helvetica", 8)
        c.drawString(36, 16, "Árvore B em Memória Secundária — AED-PG-2026")
        c.drawRightString(W - 36, 16, f"{idx+1}/{len(SLIDES)}")
        c.showPage()

    c.save()
    print("PDF:", path)

# ---------------------------------------------------------------------------
# Versão CLEAN: mesmos gráficos/tabelas/resultados, texto enxuto (a fala cobre
# os detalhes). Reaproveita todos os C_* / tabelas já gerados acima; só troca o
# conteúdo textual dos slides e os nomes de saída (APRESENTACAO_2).
# ---------------------------------------------------------------------------
SLIDES = [
    {
        "layout": "cover",
        "title": "Árvore B em Memória Secundária",
        "subtitle": "Uma classe Árvore B de ordem m operando estritamente em disco",
        "bullets": [
            "Algoritmos e Estruturas de Dados — AED-PG-2026 (5955001-3)",
            "Prof. Dr. José Augusto Baranauskas — DCM / USP",
            "C++17  ·  1º Semestre/2026",
            "Cássio Takarada  ·  Cézio Luiz Ferreira Junior",
        ],
    },
    {
        "layout": "two_col",
        "title": "Roteiro",
        "left_title": "Parte 1 — a estrutura",
        "left": [
            "Por que disco e por que Árvore B",
            "Anatomia do nó e do arquivo",
            "Busca, inserção e remoção",
            "Demonstração ao vivo",
        ],
        "right_title": "Parte 2 — a avaliação",
        "right": [
            "Experimentos e métricas",
            "Impacto da ordem m",
            "Escala e reaproveitamento",
            "Validação em 2 máquinas + conclusões",
        ],
    },
    {
        "layout": "bullets",
        "title": "O problema",
        "bullets": [
            "Árvore B de ordem m que reside em disco (arquivo)",
            "Nunca carregada inteira — 1 nó por acesso",
            "Busca, inserção e remoção balanceadas",
            "Meta extra: avaliar m, escala e 2 máquinas",
        ],
    },
    {
        "layout": "bullets",
        "title": "Contexto — memória rápida × disco lento",
        "bullets": [
            "RAM: rápida, pequena, volátil",
            "Disco: enorme, permanente, milhares de vezes mais lento",
            "Gargalo = disco, não CPU",
            "Objetivo: minimizar acessos ao disco",
        ],
    },
    {
        "layout": "bullets",
        "title": "O índice e a métrica honesta",
        "bullets": [
            "Índice = sumário que evita ler o arquivo inteiro",
            "Árvore B: índice real de SGBDs e sistemas de arquivos",
            "Métrica honesta = nº de acessos ao disco (não o relógio)",
            "O relógio engana: a escrita vai ao cache do SO",
        ],
    },
    {
        "layout": "text_image",
        "title": "Anatomia de um nó",
        "bullets": [
            "Registro de tamanho fixo, lido de uma vez",
            "Campos: n · ponteiros A[ ] · chaves K[ ]",
            "Ordem m → até m-1 chaves e m filhos",
            "1 nó = 1 acesso ao disco",
        ],
        "images": [(C_NODE, None)],
    },
    {
        "layout": "text_image",
        "title": "Layout do arquivo no disco",
        "bullets": [
            "Arquivo = vetor de registros (RRN)",
            "RRN 0 = header (root, total, free_head)",
            "offset = RRN × PAGE_SIZE → acesso direto",
            "Nós removidos → lista de livres (reuso)",
        ],
        "images": [(C_FILE, None)],
    },
    {
        "layout": "image_full",
        "title": "Decisões de implementação",
        "bullets": [
            "Ordem m = constante de compilação (M)",
            "Raiz no header · 1 nó por I/O · nunca carrega tudo",
            "Free list no próprio arquivo (reaproveitamento)",
        ],
        "images": [(G("m3_final.png"), "Árvore B (m=3, 40 chaves) exportada pelo próprio código — o número é o RRN no disco")],
    },
    {
        "layout": "bullets",
        "title": "O DiskManager — um nó por I/O",
        "bullets": [
            "Toda E/S passa por uma única classe",
            "readNode / writeNode = 1 registro de PAGE_SIZE",
            "Aqui mora o contador de acessos (métrica central)",
            "A árvore nunca é carregada inteira",
        ],
    },
    {
        "layout": "bullets",
        "title": "Reaproveitamento — a lista de livres",
        "bullets": [
            "Remoção libera nó → vira nó LIVRE",
            "freeNode encadeia na lista (free_head no header)",
            "allocNode reusa antes de crescer o arquivo",
            "Ligável (setReuse) — medido no churn",
        ],
    },
    {
        "layout": "bullets",
        "title": "Métodos — visão geral",
        "bullets": [
            "Busca m-way (mSearch)",
            "Inserção bottom-up com split (insertB)",
            "Remoção: sucessor, redistribuição e fusão (deleteB)",
            "Monitoramento: printTree · height · exportDot",
        ],
    },
    {
        "layout": "text_image",
        "title": "Busca m-way",
        "bullets": [
            "Chaves dividem o universo em faixas",
            "Desce pelo ponteiro da faixa certa",
            "Custo ≈ altura da árvore",
            "Buscar 70: apenas 2 acessos",
        ],
        "images": [(G("m3_final.png"), "A busca segue um caminho da raiz à folha")],
    },
    {
        "layout": "bullets",
        "title": "Inserção e split",
        "bullets": [
            "Bottom-up: acha a folha, insere em ordem",
            "Estourou (> m-1 chaves)? → split",
            "A mediana sobe ao pai; pode propagar e criar raiz",
            "É o que mantém a árvore balanceada",
        ],
    },
    {
        "layout": "bullets",
        "title": "Remoção — sucessor, redistribuição e fusão",
        "bullets": [
            "Nó interno: troca pelo sucessor in-order",
            "Folha pode ficar em underflow",
            "Redistribuição: empresta de um irmão",
            "Fusão: funde com irmão + chave do pai",
        ],
    },
    {
        "layout": "two_col",
        "title": "Resumo — o que o código faz",
        "left_title": "Núcleo — Árvore B em disco",
        "left": [
            "Ordem m parametrizável (M em compilação)",
            "100% em disco: 1 nó por I/O",
            "Busca · inserção (split) · remoção (merge)",
            "Persistência por RRN + header",
        ],
        "right_title": "Extras",
        "right": [
            "Contador de acessos ao disco",
            "Free list ligável (setReuse)",
            "printTree · height · exportDot (Graphviz)",
            "Benchmark CSV + automação (run_all, make_tables)",
        ],
    },
    {
        "layout": "demo",
        "title": "Demonstração — o programa rodando",
        "bullets": [
            "Compilação → menu → CRUD + métricas",
            "m=3 e m=5: a árvore fica mais rasa",
            "Captura real da execução",
        ],
        "gif": os.path.join(ASSETS, "demo_btree.gif"),
        "poster": os.path.join(ASSETS, "demo_poster.png"),
        "caption": "make M=3/M=5 → menu interativo → grafo exportado pelo programa",
    },
    {
        "layout": "image_full",
        "title": "Grafos gerados pelo programa",
        "bullets": [
            "Exportado pelo programa (Graphviz)",
            "Cada caixa = nó no disco (RRN)",
            "m=5: mais chaves por nó → árvore mais rasa",
        ],
        "images": [(os.path.join(ASSETS, "graphs_m3_m5.png"),
                    "Mesma sequência, duas ordens m — à direita a árvore é mais rasa")],
    },
    {
        "layout": "text_image",
        "title": "Os experimentos — a matriz de testes",
        "bullets": [
            "Exp.1 — impacto da ordem m (N=10⁵)",
            "Exp.2 — escala N (10³–10⁶)",
            "Exp.3 — ocupação com/sem reuso",
            "Exp.4 — CPU vs I/O",
            "Validação em 2 máquinas",
        ],
        "images": [(C_IO_N, None)],
    },
    {
        "layout": "bullets",
        "title": "Metodologia e ambiente",
        "bullets": [
            "Driver não interativo → CSV de métricas",
            "Cada configuração reconstruída do zero",
            "Modos aleatório e sequencial",
            "getrusage: CPU usuário / sistema / espera de I/O",
            "2 máquinas: Notebook × Titan (USP)",
        ],
    },
    {
        "layout": "text_image",
        "title": "Métricas utilizadas",
        "bullets": [
            "Acessos ao disco — métrica central",
            "Altura (~ log_m N)",
            "Ocupação: nós, livres e bytes",
            "Tempo: wall · CPU (user/sys) · espera de I/O",
        ],
        "images": [(C_H_M, None)],
    },
    {
        "layout": "table_image",
        "title": "Impacto da ordem m — I/O por busca (N=10⁵)",
        "table": (ORDER_HEADERS, ORDER_ROWS),
        "images": [(C_IO_M, None)],
        "footnote": "I/O por busca cai de 13,25 (m=3, altura 14) para 2,00 (m≥512, altura 2). Ganho satura por volta de m≈64-128.",
    },
    {
        "layout": "text_image",
        "title": "Impacto da ordem m — altura",
        "bullets": [
            "Altura ≈ acessos por busca",
            "m=3 → 14 níveis (fundo demais)",
            "m cresce → altura despenca em degraus",
            "m grande → altura 2",
        ],
        "images": [(C_H_M, None)],
    },
    {
        "layout": "text_image",
        "title": "Resultados por operação",
        "bullets": [
            "Busca: a mais barata (raiz → folha)",
            "Inserção: split reescreve nós",
            "Remoção: a mais cara (redistribui / funde)",
            "Todas caem com m e saturam juntas",
        ],
        "images": [(C_PHASE, None)],
    },
    {
        "layout": "text_image",
        "title": "Escala do conjunto (mil → 1 milhão)",
        "bullets": [
            "N de mil a um milhão de chaves",
            "Acessos por busca quase não mudam",
            "Custo logarítmico na base m",
            "Escalável: cresce o dado, não o custo",
        ],
        "images": [(C_IO_N, None)],
    },
    {
        "layout": "table_image",
        "title": "Reaproveitamento de nós (churn)",
        "table": (REUSE_HEADERS, REUSE_ROWS),
        "images": [(C_REUSE, None)],
        "footnote": "O reaproveitamento economiza ~27-30% do tamanho do arquivo neste workload.",
    },
    {
        "layout": "table_image",
        "title": "Dois sistemas — tempo (wall)",
        "table": (MACH_HEADERS, MACH_ROWS),
        "images": [(C_WALL_M, None)],
        "footnote": "Mesmo programa, 2 máquinas: acessos a disco idênticos; só o tempo muda — a curva tem a mesma forma.",
    },
    {
        "layout": "two_image",
        "title": "Dois sistemas — CPU e determinismo",
        "bullets": [
            "Tempo é quase todo CPU de sistema (I/O), não espera de disco",
            "Acessos idênticos: todos os pontos sobre y=x",
            "Determinística e portável",
        ],
        "images": [(C_CPU, "CPU usuário vs sistema (I/O) nas 2 máquinas"),
                   (C_DET, "Acessos a disco idênticos — pontos sobre y=x")],
    },
    {
        "layout": "text_image",
        "title": "Dois sistemas — validação cruzada",
        "bullets": [
            "Notebook × Titan (USP)",
            "Acessos a disco idênticos",
            "Só o tempo varia",
            "Comparação via compare.py",
        ],
        "images": [(C_MACH, None)] if C_MACH else [],
    },
    {
        "layout": "text_image",
        "title": "Dois sistemas — escala do conjunto",
        "bullets": [
            "N até 1 milhão (m=100), nas 2 máquinas",
            "Tempo cresce devagar (logarítmico)",
            "Mesma tendência; Titan mais rápido",
            "Prova prática de escalabilidade",
        ],
        "images": [(C_SCALE_M, None)],
    },
    {
        "layout": "two_col",
        "title": "E se usássemos union no header? (sugestão do prof.)",
        "left_title": "Como está hoje",
        "left": [
            "Header e BNode são tipos separados",
            "E/S serializa campo a campo (memcpy)",
            "Header usa 1 página inteira (RRN uniforme)",
            "Nó livre esconde o próximo em K[1]",
        ],
        "right_title": "Com union",
        "right": [
            "Página única: header e nós no mesmo tamanho fixo",
            "Lê/escreve a página inteira → menos código",
            "FreeNode.next explícito (autodocumentado)",
            "Custo: portabilidade do binário (padding/endianness)",
        ],
    },
    {
        "layout": "bullets",
        "title": "Dificuldades técnicas",
        "bullets": [
            "Indexação 1-based (RRN 0 = header)",
            "Remoção: manter o path para reparar underflow",
            "eofbit do fstream → resolvido com clear()",
            "Serialização campo a campo é propensa a erro",
        ],
    },
    {
        "layout": "two_col",
        "title": "Vantagens × Desvantagens",
        "left_title": "Vantagens",
        "left": [
            "Altura baixa → poucos acessos",
            "Sempre balanceada",
            "Ideal para disco e índices de SGBD",
            "Determinística + reuso de espaço",
        ],
        "right_title": "Desvantagens / custos",
        "right": [
            "Nós subocupados (~50% no sequencial)",
            "Remoção complexa",
            "m muito grande → mais CPU por nó",
            "Existe um m ótimo (tamanho de bloco)",
        ],
    },
    {
        "layout": "bullets",
        "title": "Aplicações práticas",
        "bullets": [
            "SGBDs: MySQL/InnoDB, PostgreSQL, Oracle",
            "Sistemas de arquivos: NTFS, APFS, ext4, Btrfs",
            "Chave-valor: SQLite, LMDB, BerkeleyDB",
            "Uso natural: índice de um banco em disco",
        ],
    },
    {
        "layout": "two_col",
        "title": "Ferramentas utilizadas",
        "left_title": "Implementação & experimentos",
        "left": [
            "C++17 (g++ -O2) · GNU Make",
            "std::fstream (binário, registros fixos)",
            "Graphviz (grafos) · Python + matplotlib",
            "python-pptx / reportlab (slides)",
        ],
        "right_title": "Processo & reprodutibilidade",
        "right": [
            "Versionamento com Git",
            "Experimentos reprodutíveis (CSV)",
            "Build parametrizável: make M=<ordem>",
            "Validação cruzada em 2 máquinas",
        ],
    },
    {
        "layout": "bullets",
        "title": "Conclusão",
        "bullets": [
            "Árvore B 100% em disco, 1 nó por I/O",
            "I/O por operação cai com m até saturar (~64-128)",
            "Free list economiza ~27-30% de espaço",
            "Determinística: acessos idênticos em 2 máquinas",
        ],
    },
    {
        "layout": "bullets",
        "title": "Referências",
        "bullets": [
            "Comer, D. (1979). The Ubiquitous B-Tree. ACM Computing Surveys.",
            "Knuth, D. The Art of Computer Programming, Vol. 3.",
            "Bayer, R. & McCreight, E. (1972). Large Ordered Indexes.",
            "Baranauskas, J. A. Slides AED-PG-2026 (Árvores B). DCM / USP.",
            "Obrigado! Perguntas?",
        ],
    },
]

if __name__ == "__main__":
    build_pptx(os.path.join(OUT, "APRESENTACAO_2.pptx"))
    build_pdf(os.path.join(OUT, "APRESENTACAO_2.pdf"))
    print("\nApresentação CLEAN gerada em:", OUT)
